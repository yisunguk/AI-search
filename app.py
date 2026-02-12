import streamlit as st
import os
import time
import uuid
from datetime import datetime, timedelta
from azure.storage.blob import BlobServiceClient, generate_blob_sas, BlobSasPermissions, generate_container_sas, ContainerSasPermissions
from azure.ai.translation.document import DocumentTranslationClient, DocumentTranslationInput, TranslationTarget
from azure.core.credentials import AzureKeyCredential
import urllib.parse
import requests
import fitz # PyMuPDF for page count
import pandas as pd
import zipfile
import io

# Search Manager Import
from search_manager import AzureSearchManager

# Chat Manager Import  
from chat_manager_v2 import AzureOpenAIChatManager
from doc_intel_manager import DocumentIntelligenceManager
import excel_manager

# Authentication imports
from utils.auth_manager import AuthManager
from modules.login_page import render_login_page
from utils.chat_history_utils import load_history, save_history, get_session_title
import extra_streamlit_components as stx

# -----------------------------
# 설정 및 비밀 관리
# -----------------------------
st.set_page_config(page_title="인텔리전트 다큐먼트", page_icon="🏗️", layout="wide")

# Custom CSS for larger tab labels and document list alignment
st.markdown("""
<style>
    /* Increase font size for tab labels */
    button[data-baseweb="tab"] {
        font-size: 20px !important;
    }
    button[data-baseweb="tab"] p {
        font-size: 20px !important;
        font-weight: 600 !important;
    }
    
    /* Document list - row alignment */
    [data-testid="stHorizontalBlock"] {
        display: flex !important;
        align-items: center !important;
        gap: 0.5rem !important;
        min-height: 42px !important;
    }
    
    /* Column layout - vertical centering */
    [data-testid="column"] {
        display: flex !important;
        flex-direction: column !important;
        justify-content: center !important;
    }
    
    /* All buttons - consistent height and sizing */
    .stButton button, .stLinkButton a {
        min-height: 38px !important;
        max-height: 38px !important;
        height: 38px !important;
        display: flex !important;
        align-items: center !important;
        justify-content: center !important;
        padding: 0.25rem 0.75rem !important;
        white-space: nowrap !important;
        font-size: 1.1rem !important;
    }
    
    /* Popover button - same height */
    button[data-testid="baseButton-header"] {
        min-height: 38px !important;
        max-height: 38px !important;
        height: 38px !important;
        display: flex !important;
        align-items: center !important;
        justify-content: center !important;
        padding: 0.25rem 0.75rem !important;
        font-size: 1.1rem !important;
    }
    
    /* Checkbox alignment */
    .stCheckbox {
        display: flex !important;
        align-items: center !important;
        justify-content: center !important;
        min-height: 38px !important;
    }
    
    /* Markdown text alignment */
    .stMarkdown {
        display: flex !important;
        align-items: center !important;
        min-height: 38px !important;
    }
    
    /* Prevent wrapping in icon columns */
    [data-testid="column"] > div {
        white-space: nowrap !important;
    }
</style>
""", unsafe_allow_html=True)

def get_secret(key):
    if key in st.secrets:
        return st.secrets[key]
    return os.environ.get(key)

# 필수 자격 증명
# 1. Storage
STORAGE_CONN_STR = get_secret("AZURE_STORAGE_CONNECTION_STRING")
CONTAINER_NAME = get_secret("AZURE_BLOB_CONTAINER_NAME") or "blob-leesunguk"

# 2. Translator
TRANSLATOR_KEY = get_secret("AZURE_TRANSLATOR_KEY")
TRANSLATOR_ENDPOINT = get_secret("AZURE_TRANSLATOR_ENDPOINT")

# 3. Search
SEARCH_ENDPOINT = get_secret("AZURE_SEARCH_ENDPOINT")
SEARCH_KEY = get_secret("AZURE_SEARCH_KEY")
SEARCH_INDEX_NAME = get_secret("AZURE_SEARCH_INDEX_NAME") or "pdf-search-index"
SEARCH_INDEXER_NAME = "pdf-indexer"
SEARCH_DATASOURCE_NAME = "blob-datasource"

# 4. Azure OpenAI
AZURE_OPENAI_ENDPOINT = get_secret("AZURE_OPENAI_ENDPOINT")
AZURE_OPENAI_KEY = get_secret("AZURE_OPENAI_KEY")
AZURE_OPENAI_DEPLOYMENT = get_secret("AZURE_OPENAI_DEPLOYMENT") or get_secret("AZURE_OPENAI_DEPLOYMENT_NAME")
AZURE_OPENAI_API_VERSION = get_secret("AZURE_OPENAI_API_VERSION")

# 5. Document Intelligence
AZURE_DOC_INTEL_ENDPOINT = get_secret("AZURE_DOC_INTEL_ENDPOINT")
AZURE_DOC_INTEL_KEY = get_secret("AZURE_DOC_INTEL_KEY")

# -----------------------------
# Azure 클라이언트 헬퍼
# -----------------------------
def get_blob_service_client():
    if not STORAGE_CONN_STR:
        st.error("Azure Storage Connection String이 설정되지 않았습니다.")
        st.stop()
    return BlobServiceClient.from_connection_string(STORAGE_CONN_STR)

def get_translation_client():
    if not TRANSLATOR_KEY or not TRANSLATOR_ENDPOINT:
        st.error("Azure Translator Key 또는 Endpoint가 설정되지 않았습니다.")
        st.stop()
    return DocumentTranslationClient(TRANSLATOR_ENDPOINT, AzureKeyCredential(TRANSLATOR_KEY))

def get_search_manager():
    if not SEARCH_ENDPOINT or not SEARCH_KEY:
        st.error("Azure Search Endpoint 또는 Key가 설정되지 않았습니다.")
        st.stop()
    return AzureSearchManager(SEARCH_ENDPOINT, SEARCH_KEY, SEARCH_INDEX_NAME)

def get_chat_manager():
    if not AZURE_OPENAI_ENDPOINT or not AZURE_OPENAI_KEY:
        st.error("Azure OpenAI Endpoint 또는 Key가 설정되지 않았습니다.")
        st.stop()
    return AzureOpenAIChatManager(
        AZURE_OPENAI_ENDPOINT, 
        AZURE_OPENAI_KEY, 
        AZURE_OPENAI_DEPLOYMENT, 
        AZURE_OPENAI_API_VERSION,
        get_search_manager(),
        STORAGE_CONN_STR,
        CONTAINER_NAME
    )

def get_doc_intel_manager():
    if not AZURE_DOC_INTEL_ENDPOINT or not AZURE_DOC_INTEL_KEY:
        st.error("Azure Document Intelligence Endpoint 또는 Key가 설정되지 않았습니다.")
        st.stop()
    return DocumentIntelligenceManager(AZURE_DOC_INTEL_ENDPOINT, AZURE_DOC_INTEL_KEY)

def generate_sas_url(blob_service_client, container_name, blob_name=None, page=None, permission="r", expiry_hours=1, content_disposition=None, no_viewer=False):
    """
    Generates a SAS URL for a blob and wraps it in a web viewer (Google Docs/Office) if applicable.
    If blob_name is None, generates a Container SAS.
    """
    try:
        account_name = blob_service_client.account_name
        
        # Handle credential types
        if hasattr(blob_service_client.credential, 'account_key'):
            account_key = blob_service_client.credential.account_key
        else:
            account_key = blob_service_client.credential['account_key']
        
        start = datetime.utcnow() - timedelta(minutes=15)
        expiry = datetime.utcnow() + timedelta(hours=expiry_hours)
        
        if blob_name:
            # Clean blob name (remove page suffixes like " (p.1)")
            import re
            clean_name = re.sub(r'\s*\(\s*p\.?\s*\d+\s*\)', '', blob_name).strip()
            
            # Determine content type
            import mimetypes
            content_type, _ = mimetypes.guess_type(clean_name)
            
            # Force PDF content type if extension matches (to ensure browser opens it)
            if clean_name.lower().endswith('.pdf'):
                content_type = "application/pdf"
                content_disposition = "inline"
            elif not content_type:
                content_type = "application/octet-stream"

            if content_disposition is None:
                content_disposition = "inline"

            sas_token = generate_blob_sas(
                account_name=account_name,
                container_name=container_name,
                blob_name=clean_name,
                account_key=account_key,
                permission=BlobSasPermissions(read=True),
                start=start,
                expiry=expiry,
                content_disposition=content_disposition,
                content_type=content_type
            )
            sas_url = f"https://{account_name}.blob.core.windows.net/{container_name}/{urllib.parse.quote(clean_name, safe='/')}?{sas_token}"
            
            if no_viewer:
                return sas_url
            
            lower_name = clean_name.lower()
            if lower_name.endswith(('.pptx', '.ppt', '.docx', '.doc', '.xlsx', '.xls')):
                encoded_sas_url = urllib.parse.quote(sas_url)
                return f"https://view.officeapps.live.com/op/view.aspx?src={encoded_sas_url}"
            elif lower_name.endswith('.pdf'):
                # Use native browser viewer (better performance/reliability than Google Viewer)
                # encoded_sas_url = urllib.parse.quote(sas_url)
                # final_url = f"https://docs.google.com/viewer?url={encoded_sas_url}"
                
                # Direct SAS URL with content_disposition=inline opens in browser PDF viewer
                final_url = sas_url
                if page:
                    final_url += f"#page={page}"
                return final_url
            else:
                return sas_url
        else:
            # Container SAS
            sas_token = generate_container_sas(
                account_name=account_name,
                container_name=container_name,
                account_key=account_key,
                permission=ContainerSasPermissions(write=True, list=True, read=True, delete=True),
                start=start,
                expiry=expiry
            )
            return f"https://{account_name}.blob.core.windows.net/{container_name}?{sas_token}"
            
    except Exception as e:
        st.error(f"SAS URL 생성 중 오류 발생 ({blob_name}): {e}")
        return "#"

# -----------------------------
# Progress Management (Resume Capability)
# -----------------------------
import json

TEMP_DIR = ".temp_analysis"
if not os.path.exists(TEMP_DIR):
    os.makedirs(TEMP_DIR)

def get_progress_file_path(safe_filename):
    return os.path.join(TEMP_DIR, f"{safe_filename}_progress.json")

def save_progress(safe_filename, page_chunks, total_pages):
    """Save intermediate analysis progress to disk"""
    try:
        filepath = get_progress_file_path(safe_filename)
        data = {
            "safe_filename": safe_filename,
            "total_pages": total_pages,
            "page_chunks": page_chunks,
            "last_updated": datetime.utcnow().isoformat()
        }
        with open(filepath, 'w', encoding='utf-8') as f:
            json.dump(data, f, ensure_ascii=False, indent=2)
        # print(f"DEBUG: Progress saved for {safe_filename} ({len(page_chunks)} chunks)")
    except Exception as e:
        print(f"Error saving progress: {e}")

def load_progress(safe_filename):
    """Load intermediate analysis progress from disk"""
    try:
        filepath = get_progress_file_path(safe_filename)
        if os.path.exists(filepath):
            with open(filepath, 'r', encoding='utf-8') as f:
                data = json.load(f)
            print(f"DEBUG: Progress loaded for {safe_filename} ({len(data.get('page_chunks', []))} chunks)")
            return data
    except Exception as e:
        print(f"Error loading progress: {e}")
    return None

def delete_progress(safe_filename):
    """Delete progress file after successful completion"""
    try:
        filepath = get_progress_file_path(safe_filename)
        if os.path.exists(filepath):
            os.remove(filepath)
            print(f"DEBUG: Progress file deleted for {safe_filename}")
    except Exception as e:
        print(f"Error deleting progress file: {e}")

# File Persistence for Resume
FILES_DIR = os.path.join(TEMP_DIR, "files")
if not os.path.exists(FILES_DIR):
    os.makedirs(FILES_DIR)

def save_uploaded_file_temp(uploaded_file, safe_filename):
    """Save uploaded file to temp dir for resume capability"""
    try:
        filepath = os.path.join(FILES_DIR, safe_filename)
        with open(filepath, "wb") as f:
            f.write(uploaded_file.getbuffer())
        return filepath
    except Exception as e:
        print(f"Error saving temp file: {e}")
        return None

def get_temp_file_path(safe_filename):
    return os.path.join(FILES_DIR, safe_filename)

def delete_temp_file(safe_filename):
    """Delete temp file after completion"""
    try:
        filepath = os.path.join(FILES_DIR, safe_filename)
        if os.path.exists(filepath):
            os.remove(filepath)
            print(f"DEBUG: Temp file deleted for {safe_filename}")
    except Exception as e:
        print(f"Error deleting temp file: {e}")

class LocalFile:
    """Mock Streamlit UploadedFile for local files"""
    def __init__(self, path, name, type="application/pdf"):
        self.path = path
        self.name = name
        self.type = type
        self.size = os.path.getsize(path)
        self._file = open(path, "rb")

    def read(self, size=-1):
        return self._file.read(size)

    def seek(self, offset, whence=0):
        return self._file.seek(offset, whence)

    def tell(self):
        return self._file.tell()
        
    def getbuffer(self):
        # Return bytes
        self.seek(0)
        return self.read()

    def close(self):
        self._file.close()

def is_drm_protected(uploaded_file):
    """
    Check if the uploaded file is DRM protected or encrypted.
    Returns True if protected, False otherwise.
    """
    try:
        file_type = uploaded_file.name.split('.')[-1].lower()
        
        # 1. PDF Check
        if file_type == 'pdf':
            try:
                # Read file stream
                bytes_data = uploaded_file.getvalue()
                with fitz.open(stream=bytes_data, filetype="pdf") as doc:
                    if doc.is_encrypted:
                        return True
            except Exception as e:
                print(f"PDF DRM Check Error: {e}")
                # If we can't open it with fitz, it might be corrupted or heavily encrypted
                return True 

        # 2. Office Files (docx, pptx, xlsx) Check
        # Modern Office files are Zip archives. If they are encrypted/DRM'd, 
        # they often become OLE CF (Compound File) binaries or non-zip streams.
        elif file_type in ['docx', 'pptx', 'xlsx']:
            try:
                bytes_data = uploaded_file.getvalue()
                # Check if it is a valid zip file
                if not zipfile.is_zipfile(io.BytesIO(bytes_data)):
                    # Not a zip -> Likely Encrypted/DRM (OLE format)
                    return True
                
                # Optional: Try to open it to be sure
                with zipfile.ZipFile(io.BytesIO(bytes_data)) as zf:
                    # Check for standard OOXML structure (e.g., [Content_Types].xml)
                    if '[Content_Types].xml' not in zf.namelist():
                        return True
            except Exception as e:
                print(f"Office DRM Check Error: {e}")
                return True # Assume protected if we can't parse structure
                
        return False
    except Exception as e:
        print(f"General DRM Check Error: {e}")
        return False



# -----------------------------
# UI 구성
# -----------------------------


# 지원 언어 목록 가져오기 (API)
@st.cache_data
def get_supported_languages():
    try:
        url = "https://api.cognitive.microsofttranslator.com/languages?api-version=3.0&scope=translation"
        # Accept-Language 헤더를 'ko'로 설정하여 언어 이름을 한국어로 받음
        headers = {"Accept-Language": "ko"}
        response = requests.get(url, headers=headers, timeout=5)
        response.raise_for_status()
        data = response.json()
        
        languages = {}
        for code, info in data['translation'].items():
            # "한국어 이름 (원어 이름)" 형식으로 표시 (예: 영어 (English))
            label = f"{info['name']} ({info['nativeName']})"
            languages[label] = code
        return languages
    except requests.exceptions.SSLError:
        # 로컬 환경(사내망) 등에서 SSL 인증서 오류 발생 시 verify=False로 재시도
        try:
            import urllib3
            urllib3.disable_warnings(urllib3.exceptions.InsecureRequestWarning)
            response = requests.get(url, headers=headers, verify=False, timeout=5)
            response.raise_for_status()
            data = response.json()
            languages = {}
            for code, info in data['translation'].items():
                label = f"{info['name']} ({info['nativeName']})"
                languages[label] = code
            return languages
        except Exception as e:
            print(f"SSL Bypass retry failed: {e}")
            # 실패 시 아래 기본 언어 제공으로 넘어감

    except Exception as e:
        print(f"언어 목록 가져오기 실패 (API): {e}")
        # UI에 에러를 표시하지 않고 콘솔에만 남김
    
    # 실패 시 기본 언어 제공 (확장된 목록)
    return {
        "한국어 (Korean)": "ko", 
        "영어 (English)": "en",
        "일본어 (Japanese)": "ja",
        "중국어 간체 (Chinese Simplified)": "zh-Hans",
        "중국어 번체 (Chinese Traditional)": "zh-Hant",
        "프랑스어 (French)": "fr",
        "독일어 (German)": "de",
        "스페인어 (Spanish)": "es",
        "러시아어 (Russian)": "ru",
        "베트남어 (Vietnamese)": "vi"
    }

LANGUAGES = get_supported_languages()

# 언어 코드별 파일 접미사 매핑 (기본적으로 대문자 코드를 사용하되, 일부 커스텀 가능)
# 여기서는 자동 생성 로직을 사용하므로 별도 딕셔너리 불필요, 
# 다만 중국어 등 특수 케이스를 위해 남겨둘 수 있음.
LANG_SUFFIX_OVERRIDE = {
    "zh-Hans": "CN",
    "zh-Hant": "TW",
}

# Initialize session state for page navigation
if "page" not in st.session_state:
    st.session_state.page = "홈"

def change_page(page_name):
    st.session_state.page = page_name

# Initialize AuthManager
auth_manager = AuthManager(STORAGE_CONN_STR)

# Initialize Cookie Manager
# Initialize Cookie Manager
cookie_manager = stx.CookieManager(key="auth_cookie_manager")

# Initialize login state
if 'is_logged_in' not in st.session_state:
    st.session_state.is_logged_in = False

# Check for existing session cookie (Auto-login)
# Check for existing session cookie (Auto-login)
if not st.session_state.is_logged_in and not st.session_state.get('just_logged_out', False):
    try:
        # Improved robust cookie retrieval (Retry mechanism)
        auth_email = None
        
        # Method 1: Direct get with retries (wait for component to sync)
        # extra_streamlit_components sometimes needs a moment to load cookies from frontend
        for i in range(5):
            auth_email = cookie_manager.get(cookie="auth_email")
            if auth_email:
                break
            time.sleep(0.1) 
            
        # Method 2: Fallback to get_all() if direct get failed
        if not auth_email:
            cookies = cookie_manager.get_all()
            if cookies and isinstance(cookies, dict):
                auth_email = cookies.get("auth_email")
        
        if auth_email:
            # Validate email exists in auth_manager
            user = auth_manager.get_user_by_email(auth_email)
            if user:
                st.session_state.is_logged_in = True
                st.session_state.user_info = user
                st.toast(f"자동 로그인되었습니다: {user.get('name')}")
    except Exception as e:
        print(f"Cookie check failed: {e}")

# Check if user is logged in
if not st.session_state.is_logged_in:
    render_login_page(auth_manager, cookie_manager)
    st.stop()

# User is logged in - get their info
user_info = st.session_state.get('user_info', {})
user_role = user_info.get('role', 'guest')
user_perms = user_info.get('permissions', [])

def get_user_folder_name(user_info):
    """Get sanitized user folder name"""
    if not user_info:
        return "guest"
    # Use name but fallback to ID if empty
    name = user_info.get('name', user_info.get('id', 'guest'))
    return name.strip()

user_folder = get_user_folder_name(user_info)

# Define role-based menu permissions (Fallback / Admin)
ALL_MENUS = ["홈", "번역하기", "파일 보관함", "물어보면 답하는 문서 AI", "도면/스펙 비교", "엑셀데이터 자동추출", "사진대지 자동작성", "작업계획 및 투입비 자동작성", "관리자 설정", "사용자 설정", "디버그 (Debug)"]
GUEST_MENUS = ["홈", "사용자 설정"]

if user_role == 'admin':
    available_menus = ALL_MENUS
else:
    # Use assigned permissions, ensuring mandatory menus are present
    available_menus = user_perms if user_perms else GUEST_MENUS
    
    # Filter out button permissions (btn:download, btn:edit, etc.) from menu list
    available_menus = [menu for menu in available_menus if not menu.startswith('btn:')]
    
    # Map old menu names to new names (Migration fix)
    available_menus = [
        "도면/스펙 비교" if menu == "도면/스펙 분석" else 
        "물어보면 답하는 문서 AI" if menu in ["검색 & AI 채팅", "문서 업로드 & AI 채팅"] else menu 
        for menu in available_menus
    ]
    # Ensure "홈" and "사용자 설정" are always available
    if "홈" not in available_menus:
        available_menus.insert(0, "홈")
    if "사용자 설정" not in available_menus:
        available_menus.append("사용자 설정")
    
    # Remove "관리자 설정" if somehow present for non-admins
    if "관리자 설정" in available_menus:
        available_menus.remove("관리자 설정")

with st.sidebar:
    # User profile
    st.markdown(f"### 👤 {user_info.get('name', 'User')}")
    st.caption(f"**{user_info.get('email', '')}**")
    st.caption(f"권한: {user_role.upper()}")
    
    # Debug: Show permissions and menus
    # st.caption(f"Perms: {user_perms}")
    # st.caption(f"Menus: {available_menus}")
    
    # --- Persistent Error Display ---
    if "drm_error_message" in st.session_state and st.session_state.drm_error_message:
        st.error(st.session_state.drm_error_message)
        # Clear it after showing
        del st.session_state.drm_error_message
    
    if st.button("🚪 로그아웃", key="logout_btn", use_container_width=True):
        st.session_state.is_logged_in = False
        st.session_state.user_info = None
        st.session_state.just_logged_out = True # Prevent immediate auto-login
        # Delete cookie
        cookie_manager.delete("auth_email")
        st.rerun()
    
    st.divider()
    
    st.header("메뉴")
    # Filter menu based on user role
    menu = st.radio("이동", available_menus, key="page")
    
    st.divider()
    
    if menu == "번역하기":
        st.header("설정")
        # 한국어를 기본값으로 찾기
        default_index = 0
        lang_labels = list(LANGUAGES.keys())
        for i, label in enumerate(lang_labels):
            if "Korean" in label or "한국어" in label:
                default_index = i
                break
                
        target_lang_label = st.selectbox("목표 언어 선택", lang_labels, index=default_index)
        target_lang_code = LANGUAGES[target_lang_label]
        st.info(f"선택된 목표 언어: {target_lang_code}")

    # 자격 증명 상태 확인
    if STORAGE_CONN_STR and TRANSLATOR_KEY and SEARCH_KEY:
        st.success("✅ Azure 자격 증명 확인됨")
    else:
        st.warning("⚠️ 일부 Azure 자격 증명이 누락되었습니다.")

# Common Header for non-Home pages - Removed to allow custom placement
# if menu != "홈":
#     st.title(menu)


if menu == "홈":
    # Use the new home_chat module with function calling support
    from home_chat import render_home_chat
    chat_manager = get_chat_manager()
    render_home_chat(chat_manager)
    
if menu == "번역하기":
    _, col_main, _ = st.columns([0.1, 0.8, 0.1])
    with col_main:
        if "translate_uploader_key" not in st.session_state:
            st.session_state.translate_uploader_key = 0

        uploaded_file = st.file_uploader("번역할 문서 업로드 (PPTX, PDF, DOCX, XLSX 등)", type=["pptx", "pdf", "docx", "xlsx"], key=f"translate_{st.session_state.translate_uploader_key}")

        # 이전 번역 결과가 있으면 표시
        if "last_translation_result" in st.session_state:
            result = st.session_state.last_translation_result
            st.success("✅ 번역이 완료되었습니다!")
            st.markdown(f"[{result['file_name']} 다운로드]({result['url']})", unsafe_allow_html=True)
            
            # 결과를 지우고 싶을 수 있으므로 닫기 버튼 제공 (선택 사항)
            if st.button("결과 닫기"):
                del st.session_state.last_translation_result
                st.rerun()

        if uploaded_file:
            if is_drm_protected(uploaded_file):
                st.session_state.drm_error_message = "⛔ DRM으로 보호된 파일(암호화된 파일)은 번역할 수 없습니다. 파일 목록에서 제거되었습니다."
                st.session_state.translate_uploader_key += 1
                st.rerun()

        if st.button("번역 시작", type="primary", disabled=not uploaded_file):
            if not uploaded_file:
                st.error("파일을 업로드해주세요.")
            else:
                with st.spinner("Azure Blob에 파일 업로드 중..."):
                    try:
                        blob_service_client = get_blob_service_client()
                        container_client = blob_service_client.get_container_client(CONTAINER_NAME)
                        
                        # 컨테이너 접근 권한 확인
                        try:
                            if not container_client.exists():
                                container_client.create_container()
                        except Exception as e:
                            if "AuthenticationFailed" in str(e):
                                st.error("🚨 인증 실패: Azure Storage Key가 올바르지 않습니다. Secrets 설정을 확인해주세요.")
                                st.stop()
                            else:
                                raise e

                        # 파일명 유니크하게 처리 (UUID 제거, 덮어쓰기 허용)
                        # file_uuid = str(uuid.uuid4())[:8] 
                        original_filename = uploaded_file.name
                        input_blob_name = f"{user_folder}/original/{original_filename}"
                        
                        # 업로드
                        blob_client = container_client.get_blob_client(input_blob_name)
                        blob_client.upload_blob(uploaded_file, overwrite=True)
                        
                        st.success("업로드 완료! 번역 요청 중...")
                        
                        # SAS 생성
                        source_url = generate_sas_url(blob_service_client, CONTAINER_NAME, input_blob_name, no_viewer=True)
                        
                        # Target URL 설정
                        target_base_url = f"https://{blob_service_client.account_name}.blob.core.windows.net/{CONTAINER_NAME}"
                        # Target URL은 컨테이너 또는 폴더 경로여야 함 (파일 경로 불가)
                        # 사용자별 translated 폴더로 설정
                        # URL 인코딩 필요
                        encoded_user_folder = urllib.parse.quote(user_folder)
                        target_output_url = f"{target_base_url}/{encoded_user_folder}/translated/?{generate_sas_url(blob_service_client, CONTAINER_NAME).split('?')[1]}"
                        
                    except Exception as e:
                        st.error(f"업로드/SAS 생성 실패: {e}")
                        st.stop()

                with st.spinner("번역 작업 요청 및 대기 중..."):
                    try:
                        client = get_translation_client()
                        
                        poller = client.begin_translation(
                            inputs=[
                                DocumentTranslationInput(
                                    source_url=source_url,
                                    storage_type="File",
                                    targets=[
                                        TranslationTarget(
                                            target_url=target_output_url,
                                            language=target_lang_code
                                        )
                                    ]
                                )
                            ]
                        )
                        
                        result = poller.result()
                        
                        for doc in result:
                            if doc.status == "Succeeded":
                                st.success(f"번역 완료! (상태: {doc.status})")
                            else:
                                st.error(f"문서 번역 실패! (상태: {doc.status})")
                                if doc.error:
                                    st.error(f"에러 코드: {doc.error.code}, 메시지: {doc.error.message}")
                        
                        # 결과 파일 찾기
                        time.sleep(2)
                        # UUID 폴더가 없으므로 translated 폴더 전체에서 해당 파일명 검색
                        output_prefix_search = f"{user_folder}/translated/"
                        output_blobs = list(container_client.list_blobs(name_starts_with=output_prefix_search))
                        
                        # 방금 번역된 파일 찾기 (파일명 매칭)
                        # Azure 번역은 원본 파일명을 유지하거나 언어 코드를 붙임
                        target_blobs = []
                        for blob in output_blobs:
                            if original_filename in blob.name:
                                target_blobs.append(blob)
                        
                        if not target_blobs:
                            st.warning(f"결과 파일을 찾는 중입니다... (경로: {output_prefix_search})")
                            # Fallback: list all to debug
                            # all_output = list(container_client.list_blobs(name_starts_with=output_prefix_search))
                            # debug_msg = "\n".join([b.name for b in all_output[:10]])
                            # st.error(f"결과 파일을 찾을 수 없습니다.\n현재 폴더 파일 목록:\n{debug_msg}")
                        else:
                            st.subheader("다운로드")
                            for blob in target_blobs:
                                blob_name = blob.name
                                file_name = blob_name.split("/")[-1]
                                
                                # 파일명에 언어 접미사 추가 (Rename)
                                suffix = LANG_SUFFIX_OVERRIDE.get(target_lang_code, target_lang_code.upper())
                                name_part, ext_part = os.path.splitext(file_name)
                                
                                # 이미 접미사가 있는지 확인 (혹시 모를 중복 방지)
                                if not name_part.endswith(f"_{suffix}"):
                                    new_file_name = f"{name_part}_{suffix}{ext_part}"
                                    new_blob_name = f"{user_folder}/translated/{new_file_name}"
                                    
                                    try:
                                        # Rename: Copy to new name -> Delete old
                                        source_blob = container_client.get_blob_client(blob_name)
                                        dest_blob = container_client.get_blob_client(new_blob_name)
                                        
                                        source_sas = generate_sas_url(blob_service_client, CONTAINER_NAME, blob_name)
                                        dest_blob.start_copy_from_url(source_sas)
                                        
                                        # Wait for copy
                                        for _ in range(10):
                                            props = dest_blob.get_blob_properties()
                                            if props.copy.status == "success":
                                                break
                                            time.sleep(0.2)
                                            
                                        source_blob.delete_blob()
                                        
                                        # Update variables for download link
                                        blob_name = new_blob_name
                                        file_name = new_file_name
                                        st.toast(f"파일명 변경됨: {file_name}")
                                        
                                    except Exception as e:
                                        st.warning(f"파일명 변경 실패 (기본 이름으로 유지): {e}")

                                # PPTX 폰트 변경 (Times New Roman)
                                if file_name.lower().endswith(".pptx"):
                                    try:
                                        from pptx import Presentation
                                        
                                        # 임시 파일로 다운로드
                                        temp_pptx = f"temp_{original_filename}"
                                        blob_client_temp = container_client.get_blob_client(blob_name)
                                        with open(temp_pptx, "wb") as f:
                                            data = blob_client_temp.download_blob().readall()
                                            f.write(data)
                                        
                                        # 폰트 변경 로직
                                        prs = Presentation(temp_pptx)
                                        font_name = "Times New Roman"
                                        
                                        def change_font(shapes):
                                            for shape in shapes:
                                                if shape.has_text_frame:
                                                    for paragraph in shape.text_frame.paragraphs:
                                                        for run in paragraph.runs:
                                                            run.font.name = font_name
                                                
                                                if shape.has_table:
                                                    for row in shape.table.rows:
                                                        for cell in row.cells:
                                                            if cell.text_frame:
                                                                for paragraph in cell.text_frame.paragraphs:
                                                                    for run in paragraph.runs:
                                                                        run.font.name = font_name
                                                
                                                if shape.shape_type == 6: # Group
                                                    change_font(shape.shapes)

                                        for slide in prs.slides:
                                            change_font(slide.shapes)
                                        
                                        prs.save(temp_pptx)
                                        
                                        # 다시 업로드 (덮어쓰기)
                                        with open(temp_pptx, "rb") as f:
                                            blob_client_temp.upload_blob(f, overwrite=True)
                                        
                                        os.remove(temp_pptx)
                                        st.toast("PPTX 폰트 변경 완료 (Times New Roman)")
                                        
                                    except Exception as e:
                                        st.warning(f"PPTX 폰트 변경 실패: {e}")

                                download_sas = generate_sas_url(blob_service_client, CONTAINER_NAME, blob_name, no_viewer=True, content_disposition="attachment")
                                st.markdown(f"[{file_name} 다운로드]({download_sas})", unsafe_allow_html=True)
                                
                                # 결과 세션에 저장
                                st.session_state.last_translation_result = {
                                    "file_name": file_name,
                                    "url": download_sas
                                }
                                
                        # 성공적으로 완료되면 업로더 초기화 (키 변경)
                        st.session_state.translate_uploader_key += 1
                        time.sleep(1) # 잠시 대기
                        st.rerun()
                                
                    except Exception as e:
                        st.error(f"번역 요청 중 오류 발생: {e}")

elif menu == "파일 보관함":
    _, col_main, _ = st.columns([0.1, 0.8, 0.1])
    with col_main:
        # st.subheader("📂 클라우드 파일 보관함") - Removed to avoid duplication
        
        st.divider()
        
        if st.button("🔄 목록 새로고침"):
            st.rerun()
            
        try:
            blob_service_client = get_blob_service_client()
            container_client = blob_service_client.get_container_client(CONTAINER_NAME)
            
            # 탭으로 Input/Output 구분
            tab1, tab2 = st.tabs(["원본 문서 (Input)", "번역된 문서 (Output)"])
            
            def render_file_list(prefixes, tab_name):
                all_blobs = []
                for prefix in prefixes:
                    blobs = list(container_client.list_blobs(name_starts_with=prefix))
                    all_blobs.extend(blobs)
                
                # 중복 제거 (혹시 모를 경우 대비)
                unique_blobs = {b.name: b for b in all_blobs}.values()
                blobs = list(unique_blobs)
                blobs.sort(key=lambda x: x.creation_time, reverse=True)
                
                if not blobs:
                    st.info(f"{tab_name}에 파일이 없습니다.")
                    return

                for i, blob in enumerate(blobs):
                    file_name = blob.name.split("/")[-1]
                    creation_time = blob.creation_time.strftime('%Y-%m-%d %H:%M')
                    
                    # 폴더 경로 표시 (관리자 편의)
                    folder_path = "/".join(blob.name.split("/")[:-1])
                    
                    with st.container():
                        col1, col2, col3 = st.columns([6, 2, 2])
                        
                        with col1:
                            sas_url = generate_sas_url(blob_service_client, CONTAINER_NAME, blob.name)
                            st.markdown(f"**[{file_name}]({sas_url})**")
                            st.caption(f"📂 {folder_path} | 📅 {creation_time} | 📦 {blob.size / 1024:.1f} KB")
                        
                        with col2:
                            # 수정 (이름 변경)
                            with st.popover("수정"):
                                new_name = st.text_input("새 파일명", value=file_name, key=f"rename_{i}_{blob.name}")
                                if st.button("이름 변경", key=f"btn_rename_{i}_{blob.name}"):
                                    try:
                                        # 새 경로 생성 (기존 폴더 구조 유지)
                                        path_parts = blob.name.split("/")
                                        folder = "/".join(path_parts[:-1])
                                        new_blob_name = f"{folder}/{new_name}"
                                        
                                        # 복사 (Rename은 Copy + Delete)
                                        source_blob = container_client.get_blob_client(blob.name)
                                        dest_blob = container_client.get_blob_client(new_blob_name)
                                        
                                        # SAS URL for Copy Source
                                        source_sas = generate_sas_url(blob_service_client, CONTAINER_NAME, blob.name)
                                        
                                        dest_blob.start_copy_from_url(source_sas)
                                        
                                        # 복사 완료 대기 (간단한 폴링)
                                        for _ in range(10):
                                            props = dest_blob.get_blob_properties()
                                            if props.copy.status == "success":
                                                break
                                            time.sleep(0.5)
                                        
                                        # 원본 삭제
                                        source_blob.delete_blob()
                                        st.success("이름 변경 완료!")
                                        time.sleep(1)
                                        st.rerun()
                                    except Exception as e:
                                        st.error(f"이름 변경 실패: {e}")

                        with col3:
                            # 삭제
                            if st.button("삭제", key=f"del_{prefix}_{i}", type="secondary"):
                                try:
                                    container_client.delete_blob(blob.name)
                                    st.success("삭제되었습니다.")
                                    time.sleep(1)
                                    st.rerun()
                                except Exception as e:
                                    st.error(f"삭제 실패: {e}")
                        
                        st.divider()

            with tab1:
                input_prefixes = [f"{user_folder}/documents/"]
                if user_role == 'admin':
                    input_prefixes.extend(["input/", "gulflng/"])
                render_file_list(input_prefixes, "내 문서 (Documents)")
                
            with tab2:
                output_prefixes = [f"{user_folder}/translated/"]
                if user_role == 'admin':
                    output_prefixes.extend(["output/"])
                render_file_list(output_prefixes, "번역된 문서")
                    
        except Exception as e:
            st.error(f"파일 목록을 불러오는 중 오류 발생: {e}")

elif menu == "물어보면 답하는 문서 AI":
    from utils.chat_history_utils import load_history, save_history, get_session_title
    SEARCH_HISTORY_FILE = "search_chat_history.json"

    # Initialize Session State for Search History
    if "search_chat_history_data" not in st.session_state:
        st.session_state.search_chat_history_data = load_history(SEARCH_HISTORY_FILE)
    
    if "current_search_session_id" not in st.session_state:
        new_id = str(uuid.uuid4())
        st.session_state.current_search_session_id = new_id
        st.session_state.search_chat_history_data[new_id] = {
            "title": "새로운 대화",
            "timestamp": datetime.now().isoformat(),
            "messages": []
        }
        st.session_state.chat_messages = [] # This maps to the current session messages

    # Layout: Spacer L (25%) | Main Content (50%) | Spacer R (10%) | History Sidebar (15%)
    col_spacer_l, col_main, col_spacer_r, col_history = st.columns([0.25, 0.5, 0.1, 0.15])
    
    # Custom CSS for Sidebar Styling (Same as Home)
    st.markdown("""
    <style>
    /* Target the fourth column (History Sidebar) */
    [data-testid="stHorizontalBlock"] > [data-testid="column"]:nth-of-type(4) {
        background-color: #1E1E1E;
        border-left: 1px solid #333;
        padding: 1rem;
        border-radius: 10px;
    }
    [data-testid="column"]:nth-of-type(4) button {
        text-align: left;
    }
    </style>
    """, unsafe_allow_html=True)

    # --- Right Sidebar (History) ---
    with col_history:
        st.markdown("### 채팅 기록")
        
        if st.button("➕ 새 채팅", key="new_search_chat", use_container_width=True):
            new_id = str(uuid.uuid4())
            st.session_state.current_search_session_id = new_id
            st.session_state.search_chat_history_data[new_id] = {
                "title": "새로운 대화",
                "timestamp": datetime.now().isoformat(),
                "messages": []
            }
            st.session_state.chat_messages = []
            st.rerun()
            
        st.markdown("---")
        
        sorted_sessions = sorted(
            st.session_state.search_chat_history_data.items(),
            key=lambda x: x[1].get("timestamp", ""),
            reverse=True
        )
        
        for session_id, session_data in sorted_sessions:
            title = session_data.get("title", "대화")
            if session_id == st.session_state.current_search_session_id:
                if st.button(f"📂 {title}", key=f"search_hist_{session_id}", use_container_width=True, type="primary"):
                    pass
            else:
                if st.button(f"📄 {title}", key=f"search_hist_{session_id}", use_container_width=True):
                    st.session_state.current_search_session_id = session_id
                    st.session_state.chat_messages = session_data.get("messages", [])
                    st.rerun()
        
        if st.button("🗑️ 기록 삭제", key="del_search_hist", use_container_width=True):
            st.session_state.search_chat_history_data = {}
            save_history(SEARCH_HISTORY_FILE, {})
            new_id = str(uuid.uuid4())
            st.session_state.current_search_session_id = new_id
            st.session_state.search_chat_history_data[new_id] = {
                "title": "새로운 대화",
                "timestamp": datetime.now().isoformat(),
                "messages": []
            }
            st.session_state.chat_messages = []
            st.rerun()

    with col_main:
        st.title("물어보면 답하는 문서 AI")
        # Tabs for Search and Chat to preserve state
        tab1, tab2, tab3 = st.tabs(["📤 문서 등록", "🔎 키워드 검색", "🤖 AI 질의응답"])
        
        with tab1:
            # File Uploader (Simplified)
            if "doc_search_uploader_key" not in st.session_state:
                st.session_state.doc_search_uploader_key = 0
                
            doc_upload = st.file_uploader("문서를 등록하면 검색과 질의응답이 가능합니다.", type=['pdf', 'docx', 'txt', 'pptx'], key=f"doc_search_upload_{st.session_state.doc_search_uploader_key}")
            
            if doc_upload:
                if is_drm_protected(doc_upload):
                    st.session_state.drm_error_message = "⛔ DRM으로 보호된 파일(암호화된 파일)은 업로드할 수 없습니다. 보안을 해제한 후 다시 시도해주세요."
                    st.session_state.doc_search_uploader_key += 1
                    st.rerun()

            if doc_upload and st.button("업로드", key="btn_doc_upload"):
                try:
                    blob_service_client = get_blob_service_client()
                    container_client = blob_service_client.get_container_client(CONTAINER_NAME)
                    
                    # Upload to {user_folder}/documents/ (Flat structure)
                    blob_name = f"{user_folder}/documents/{doc_upload.name}"
                    blob_client = container_client.get_blob_client(blob_name)
                    blob_client.upload_blob(doc_upload, overwrite=True)
                    st.success(f"'{doc_upload.name}' 업로드 완료! (인덱싱에 시간이 걸릴 수 있습니다)")
                except Exception as e:
                    st.error(f"업로드 실패: {e}")
            
            st.divider()
            
            # Indexed Document List
            st.markdown("### 🗂️ 등록 문서 목록")
            
            try:
                search_manager = get_search_manager()
                
                # Construct prefix URL for filtering
                account_name = get_blob_service_client().account_name
                encoded_user_folder = urllib.parse.quote(user_folder)
                prefix_url = f"https://{account_name}.blob.core.windows.net/{CONTAINER_NAME}/{encoded_user_folder}/"
                
                # Filter logic
                if user_role == 'admin':
                    filter_expr = None
                else:
                    # Workaround: Use range query instead of startswith if startswith is not supported
                    # prefix_url ends with '/' (ASCII 47). Next char is '0' (ASCII 48).
                    # So we want path >= prefix_url AND path < prefix_url_with_next_char
                    # Actually, let's just use the next char logic safely.
                    # Or just try startswith again with debug? No, let's try the range.
                    # prefix_url = .../
                    # upper_bound = ...0
                    upper_bound = prefix_url[:-1] + '0'
                    filter_expr = f"metadata_storage_path ge '{prefix_url}' and metadata_storage_path lt '{upper_bound}'"
                
                # Debug
                # st.write(f"Debug Filter: {filter_expr}")
                
                # Search all documents (*)
                results = search_manager.search("*", filter_expr=filter_expr, top=1000)
                
                # Filter out .json files first
                filtered_results = []
                for res in results:
                    file_name = res.get('metadata_storage_name', 'Unknown')
                    if not file_name.lower().endswith('.json'):
                        filtered_results.append(res)
                
                if not filtered_results:
                    st.info("인덱싱된 문서가 없습니다.")
                else:
                    st.write(f"총 {len(filtered_results)}개 문서가 등록되어 있습니다. (검색과 질의가 가능합니다).")
                    
                    # Display as a table
                    doc_data = []
                    for res in filtered_results:
                        file_name = res.get('metadata_storage_name', 'Unknown')
                        size = res.get('metadata_storage_size', 0)
                        last_modified = res.get('metadata_storage_last_modified', '')
                        path = res.get('metadata_storage_path', '')
                        
                        # Convert size to MB
                        size_mb = f"{int(size) / (1024 * 1024):.2f} MB"
                        
                        # Format date
                        try:
                            dt = datetime.fromisoformat(last_modified.replace('Z', '+00:00'))
                            date_str = dt.strftime("%Y-%m-%d %H:%M")
                        except:
                            date_str = last_modified
                            
                        doc_data.append({
                            "Name": file_name,
                            "Size": size_mb,
                            "Last Modified": date_str,
                            "path": path # Hidden for logic
                        })
                    
                    # Use dataframe for better display
                    import pandas as pd
                    df = pd.DataFrame(doc_data)
                    
                    # Display table with selection (optional, maybe just list)
                    # For now, just a clean dataframe display
                    st.dataframe(
                        df[["Name", "Size", "Last Modified"]],
                        use_container_width=True,
                        hide_index=True
                    )
            except Exception as e:
                st.error(f"문서 목록을 불러오는 중 오류가 발생했습니다: {e}")
                if 'filter_expr' in locals():
                    st.code(filter_expr)
        
        with tab2:
            # -----------------------------
            # 검색 옵션
            # -----------------------------
            with st.expander("⚙️ 고급 검색 옵션 (RAG 설정)", expanded=False):
                c1, c2 = st.columns(2)
                with c1:
                    search_use_semantic = st.checkbox("시맨틱 랭커 사용", value=True, key="search_use_semantic", help="의미 기반 검색을 사용하여 정확도를 높입니다.")
                with c2:
                    search_mode_opt = st.radio("검색 모드", ["all (AND)", "any (OR)"], index=1, horizontal=True, key="search_mode_opt", help="any: 키워드 중 하나라도 포함되면 검색 (추천)")
                    search_mode = "all" if "all" in search_mode_opt else "any"

            # Display Chat History (Shared with AI Chat)
            for msg in st.session_state.chat_messages:
                with st.chat_message(msg["role"]):
                    if msg["role"] == "user":
                        st.markdown(msg["content"])
                    else:
                        # Assistant message (Results)
                        if "results" in msg:
                            results = msg["results"]
                            if not results:
                                st.info("검색 결과가 없습니다.")
                            else:
                                st.success(f"총 {len(results)}개의 문서를 찾았습니다.")
                                for result in results:
                                    with st.container():
                                        file_name = result.get('metadata_storage_name', 'Unknown File')
                                        path = result.get('metadata_storage_path', '')
                                        
                                        # Highlights
                                        highlights = result.get('@search.highlights')
                                        if highlights:
                                            snippets = []
                                            if 'content' in highlights:
                                                snippets.extend(highlights['content'])
                                            if 'content_exact' in highlights:
                                                snippets.extend(highlights['content_exact'])
                                            unique_snippets = list(set(snippets))[:3]
                                            content_snippet = " ... ".join(unique_snippets)
                                        else:
                                            content_snippet = result.get('content', '')[:300] + "..."
                                        
                                        # Text Cleaning Logic (Restored & Improved)
                                        import re
                                        def clean_text(text):
                                            # 1. Escape Markdown special characters except HTML tags we want to keep (like <mark>)
                                            text = text.replace('~', '\\~')
                                            
                                            # 2. Handle HTML Tables (convert to Markdown-ish for display)
                                            text = re.sub(r'</td>', ' | ', text, flags=re.IGNORECASE)
                                            text = re.sub(r'</th>', ' | ', text, flags=re.IGNORECASE)
                                            text = re.sub(r'</tr>', '\n', text, flags=re.IGNORECASE)
                                            
                                            # 3. If text contains pipes (|), it might be a table. Preserve structure.
                                            if "|" in text:
                                                # Remove other HTML tags except <mark>
                                                text = re.sub(r'<(?!/?mark\b)[^>]+>', '', text)
                                                text = re.sub(r'^\s*(\|[\s\|]*)+\s*$', '', text, flags=re.MULTILINE)
                                                text = re.sub(r'\n\s*\n', '\n', text)
                                                return text.strip()

                                            # 4. Remove other HTML tags except <mark>
                                            text = re.sub(r'<(?!/?mark\b)[^>]+>', '', text)
                                            
                                            # 5. Replace single newlines with space
                                            cleaned = re.sub(r'(?<!\.)\n(?!\n)', ' ', text)
                                            cleaned = re.sub(r' +', ' ', cleaned)
                                            return cleaned.strip()
                                            
                                        st.markdown(f"### 📄 {file_name}")
                                        st.markdown(f"> {clean_text(content_snippet)}", unsafe_allow_html=True)
                                        
                                        # Generate SAS link
                                        try:
                                            blob_service_client = get_blob_service_client()
                                            from urllib.parse import unquote
                                            
                                            if "https://direct_fetch/" in path:
                                                blob_path = unquote(path.replace("https://direct_fetch/", "").split('#')[0])
                                            elif CONTAINER_NAME in path:
                                                blob_path = unquote(path.split(f"/{CONTAINER_NAME}/")[1].split('#')[0])
                                            else:
                                                blob_path = path
                                            
                                            import mimetypes
                                            content_type, _ = mimetypes.guess_type(file_name)
                                            
                                            sas_token = generate_blob_sas(
                                                account_name=blob_service_client.account_name,
                                                container_name=CONTAINER_NAME,
                                                blob_name=blob_path,
                                                account_key=blob_service_client.credential.account_key,
                                                permission=BlobSasPermissions(read=True),
                                                expiry=datetime.utcnow() + timedelta(hours=1),
                                                content_disposition="inline",
                                                content_type=content_type
                                            )
                                            sas_url = f"https://{blob_service_client.account_name}.blob.core.windows.net/{CONTAINER_NAME}/{urllib.parse.quote(blob_path)}?{sas_token}"
                                            
                                            lower_name = file_name.lower()
                                            if lower_name.endswith(('.pptx', '.ppt', '.docx', '.doc', '.xlsx', '.xls')):
                                                final_url = f"https://view.officeapps.live.com/op/view.aspx?src={urllib.parse.quote(sas_url)}"
                                                link_text = "📄 웹에서 보기 (Office Viewer)"
                                            elif lower_name.endswith('.pdf'):
                                                final_url = f"https://docs.google.com/viewer?url={urllib.parse.quote(sas_url)}"
                                                link_text = "📄 웹에서 보기 (PDF Viewer)"
                                            else:
                                                final_url = sas_url
                                                link_text = "📄 문서 열기 (새 탭)"
                                                
                                            st.markdown(f'<a href="{final_url}" target="_blank">{link_text}</a>', unsafe_allow_html=True)
                                        except Exception as e:
                                            st.caption(f"링크 생성 실패: {e}")
                                        st.divider()
                        else:
                            st.markdown(msg["content"])

            # Chat Input for Search
            if query := st.chat_input("검색할 키워드를 입력하세요...", key="keyword_chat_input"):
                st.session_state.chat_messages.append({"role": "user", "content": query})
                
                with st.spinner("검색 중..."):
                    try:
                        search_manager = get_search_manager()
                        account_name = get_blob_service_client().account_name
                        encoded_user_folder = urllib.parse.quote(user_folder)
                        prefix_url = f"https://{account_name}.blob.core.windows.net/{CONTAINER_NAME}/{encoded_user_folder}/"
                        
                        if user_role == 'admin':
                            filter_expr = None
                        else:
                            upper_bound = prefix_url[:-1] + '0'
                            filter_expr = f"metadata_storage_path ge '{prefix_url}' and metadata_storage_path lt '{upper_bound}'"
                        
                        results = search_manager.search(query, filter_expr=filter_expr, use_semantic_ranker=search_use_semantic, search_mode=search_mode)
                        
                        # Filter out .json files
                        filtered_results = [res for res in results if not res.get('metadata_storage_name', '').lower().endswith('.json')]
                        
                        st.session_state.chat_messages.append({
                            "role": "assistant",
                            "content": f"'{query}'에 대한 검색 결과입니다.",
                            "results": filtered_results
                        })
                        
                        # --- Auto-Save History ---
                        current_id = st.session_state.current_search_session_id
                        current_title = st.session_state.search_chat_history_data[current_id]["title"]
                        if current_title == "새로운 대화" and len(st.session_state.chat_messages) > 0:
                            new_title = get_session_title(st.session_state.chat_messages)
                            st.session_state.search_chat_history_data[current_id]["title"] = new_title
                        
                        st.session_state.search_chat_history_data[current_id]["messages"] = st.session_state.chat_messages
                        st.session_state.search_chat_history_data[current_id]["timestamp"] = datetime.now().isoformat()
                        save_history(SEARCH_HISTORY_FILE, st.session_state.search_chat_history_data)
                        st.rerun()
                    except Exception as e:
                        st.error(f"검색 중 오류 발생: {e}")

        with tab3:
            # -----------------------------
            # 검색 옵션 (Chat Tab) - Moved to Top
            # -----------------------------
            with st.expander("⚙️ 고급 검색 옵션 (RAG 설정)", expanded=False):
                c1, c2 = st.columns(2)
                with c1:
                    chat_use_semantic = st.checkbox("시맨틱 랭커 사용", value=True, key="chat_use_semantic", help="의미 기반 검색을 사용하여 정확도를 높입니다.")
                with c2:
                    chat_search_mode_opt = st.radio("검색 모드", ["all (AND)", "any (OR)"], index=1, horizontal=True, key="chat_search_mode", help="any: 키워드 중 하나라도 포함되면 검색 (추천)")
                    chat_search_mode = "all" if "all" in chat_search_mode_opt else "any"
            
            # Initialize chat history in session state
            if "chat_messages" not in st.session_state:
                st.session_state.chat_messages = []
            
            # Display chat messages
            for message in st.session_state.chat_messages:
                with st.chat_message(message["role"]):
                    st.markdown(message["content"])
                    
                    # Display citations if present
                    if "citations" in message and message["citations"]:
                        st.markdown("---")
                        st.caption("📚 **참조 문서:**")
                        for i, citation in enumerate(message["citations"], 1):
                            filepath = citation.get('filepath', 'Unknown')
                            # Use pre-generated final_url if available, otherwise generate one
                            display_url = citation.get('final_url')
                            if not display_url:
                                try:
                                    blob_service_client = get_blob_service_client()
                                    display_url = generate_sas_url(
                                        blob_service_client, 
                                        CONTAINER_NAME, 
                                        filepath, 
                                        page=citation.get('page')
                                    )
                                except:
                                    display_url = "#"
                            
                            st.markdown(f"{i}. [{filepath}]({display_url})")
            
            # -----------------------------
            # 검색 옵션 (Chat Tab) - Moved to Top
            # -----------------------------
            # st.write("")
            # with st.expander("⚙️ 고급 검색 옵션 (RAG 설정)", expanded=False):
            #     c1, c2 = st.columns(2)
            #     with c1:
            #         chat_use_semantic = st.checkbox("시맨틱 랭커 사용", value=False, key="chat_use_semantic", help="의미 기반 검색을 사용하여 정확도를 높입니다.")
            #     with c2:
            #         chat_search_mode_opt = st.radio("검색 모드", ["all (AND)", "any (OR)"], index=1, horizontal=True, key="chat_search_mode", help="any: 키워드 중 하나라도 포함되면 검색 (추천)")
            #         chat_search_mode = "all" if "all" in chat_search_mode_opt else "any"

            # Chat input
            if prompt := st.chat_input("질문을 입력하세요 (예: 10-P-101A의 사양은 무엇인가요?)", key="search_chat_input"):
                # Add user message to chat history
                st.session_state.chat_messages.append({"role": "user", "content": prompt})
                
                # Display user message
                with st.chat_message("user"):
                    st.markdown(prompt)
                
                # Get AI response
                with st.chat_message("assistant"):
                    with st.spinner("답변 생성 중..."):
                        try:
                            chat_manager = get_chat_manager()
                            
                            # Prepare conversation history (exclude citations from history)
                            conversation_history = [
                                {"role": msg["role"], "content": msg["content"]}
                                for msg in st.session_state.chat_messages[:-1]  # Exclude the just-added user message
                            ]
                            
                            # Pass the selected search options to the chat manager
                            response_text, citations, context, final_filter, search_results = chat_manager.get_chat_response(
                                prompt, 
                                conversation_history, 
                                search_mode=chat_search_mode, 
                                use_semantic_ranker=chat_use_semantic,
                                filter_expr=None,
                                user_folder=user_folder, # Pass Name-based folder (matches Blob/Index path)
                                is_admin=(user_role == 'admin')
                            )
                            
                            # ---------------------------------------------------------
                            # CRITICAL: Linkify Inline Citations & Escape Tildes
                            # ---------------------------------------------------------
                            
                            # 1. Pre-generate Web Viewer URLs for all citations
                            citation_links = {}
                            processed_citations = []
                            
                            if citations:
                                for cit in citations:
                                    filepath = cit.get('filepath', 'Unknown')
                                    # CRITICAL: Clean filepath from page suffixes like " (p.1)" or " (p.1) (p.1)"
                                    import re
                                    clean_filepath = re.sub(r'\s*\(\s*p\.?\s*\d+\s*\)', '', filepath).strip()
                                    
                                    page = cit.get('page')
                                    url = cit.get('url', '')
                                    
                                    # Generate Web Viewer URL
                                    try:
                                        blob_service_client = get_blob_service_client()
                                        final_url = generate_sas_url(
                                            blob_service_client, 
                                            CONTAINER_NAME, 
                                            clean_filepath, 
                                            page=page
                                        )
                                    except Exception as e:
                                        st.error(f"URL 생성 실패 ({clean_filepath}): {e}")
                                        final_url = "#"
                                    
                                    cit['final_url'] = final_url
                                    processed_citations.append(cit)
                                    
                                    filename = os.path.basename(filepath)
                                    if page:
                                        citation_links[(filename, str(page))] = final_url
                            
                            # 2. Replace text citations with Markdown links (Support both [] and ())
                            if response_text:
                                import re
                                # Pattern to match [filename: p.1] or (filename: p.1)
                                # Improved: Allow parentheses in filenames (common in EPC drawings)
                                # CRITICAL FIX: Exclude pipe (|) to prevent crossing table boundaries
                                pattern = r'[\[\(]([^\[\]|]+?:\s*p\.?\s*(\d+))[\]\)]'
                                
                                def replace_citation(match):
                                    content = match.group(1).strip()
                                    
                                    # Remove "문서명:" prefix if present (Common in Korean LLM outputs)
                                    content = re.sub(r'^문서명\s*:\s*', '', content)
                                    
                                    # Split by last colon to separate filename and page
                                    if ':' in content:
                                        fname = content.rsplit(':', 1)[0].strip()
                                        p_num = match.group(2)
                                    else:
                                        return match.group(0)
                                        
                                    original_text = match.group(0)
                                    
                                    # Try to find matching citation with fuzzy logic
                                    target_url = None
                                    
                                    # 1. Clean LLM filename for comparison
                                    clean_llm = re.sub(r'\.pdf$', '', fname.lower().strip())
                                    
                                    for (k_fname, k_page), url in citation_links.items():
                                        # 2. Clean known filename (remove .pdf and (p.N) suffixes)
                                        clean_known = re.sub(r'\.pdf$', '', k_fname.lower().strip())
                                        clean_known = re.sub(r'\s*\(\s*p\.?\s*\d+\s*\)', '', clean_known).strip()
                                        
                                        # CRITICAL FIX: Skip empty filenames to prevent false positive matches
                                        if not clean_known:
                                            continue

                                        # 3. Match page and filename (fuzzy)
                                        if str(k_page) == str(p_num):
                                            # Exact match or containment (handling "문서명: " residue if regex failed)
                                            if clean_llm == clean_known or clean_llm in clean_known or clean_known in clean_llm:
                                                target_url = url
                                                # CRITICAL FIX: Capture the matched filename to replace text
                                                matched_filename = k_fname
                                                break
                                    
                                    if target_url:
                                        # Use Markdown link for table compatibility
                                        # Escape parentheses in URL to avoid breaking Markdown link
                                        safe_url = target_url.replace('(', '%28').replace(')', '%29')
                                        
                                        # CRITICAL FIX: Replace original text (e.g. "Same Document") with actual filename
                                        # Reconstruct text: (Filename: p.N)
                                        if matched_filename:
                                            new_text = f"({matched_filename}: p.{p_num})"
                                            return f"**[{new_text}]({safe_url})**"
                                        
                                        return f"**[{original_text}]({safe_url})**"
                                    
                                    return original_text

                                # DEBUG: Show raw response before linkification
                                st.code(response_text, language="markdown")
                                
                                response_text = re.sub(pattern, replace_citation, response_text)

                                # 3. Escape tildes
                                response_text = response_text.replace('~', '\\~')
                        
                            # Display response
                            st.markdown(response_text, unsafe_allow_html=True)
                            
                            # Display citations
                            if processed_citations:
                                st.markdown("---")
                                st.caption("📚 **참조 문서:**")
                                for i, citation in enumerate(processed_citations, 1):
                                    filepath = citation.get('filepath', 'Unknown')
                                    filename = os.path.basename(filepath)
                                    display_url = citation.get('final_url', '#')
                                    
                                    link_text = "문서 보기"
                                    if "docs.google.com" in display_url: link_text = "PDF Viewer"
                                    elif "view.officeapps" in display_url: link_text = "Office Viewer"
                                    
                                    st.markdown(f"{i}. [{filename}]({display_url}) - {link_text}")
                            
                            # Debug: Show Citation Links (Hidden by default)
                            # with st.expander("🔍 링크 디버깅 (Debug Links)", expanded=False):
                            #     st.write("Citation Links Keys:", list(citation_links.keys()))
                            #     st.write("Processed Citations:", processed_citations)
                            
                            # Add assistant response to chat history
                            st.session_state.chat_messages.append({
                                "role": "assistant",
                                "content": response_text,
                                "citations": citations,
                                "context": context,
                                "debug_filter": final_filter
                            })
                            
                            # Debug: Show Context
                            with st.expander("🔍 검색된 컨텍스트 확인 (Debug Context)", expanded=False):
                                if final_filter:
                                    st.caption(f"**OData Filter:** `{final_filter}`")
                                st.text_area("LLM에게 전달된 원문 데이터", value=context, height=300)
                            
                            # --- Auto-Save History ---
                            current_id = st.session_state.current_search_session_id
                            current_title = st.session_state.search_chat_history_data[current_id]["title"]
                            if current_title == "새로운 대화" and len(st.session_state.chat_messages) > 0:
                                new_title = get_session_title(st.session_state.chat_messages)
                                st.session_state.search_chat_history_data[current_id]["title"] = new_title
                            
                            st.session_state.search_chat_history_data[current_id]["messages"] = st.session_state.chat_messages
                            st.session_state.search_chat_history_data[current_id]["timestamp"] = datetime.now().isoformat()
                            save_history(SEARCH_HISTORY_FILE, st.session_state.search_chat_history_data)
                            st.rerun()

                        except Exception as e:
                            st.error(f"오류가 발생했습니다: {str(e)}")
        

elif menu == "도면/스펙 비교":
    DRAWING_HISTORY_FILE = "drawing_chat_history.json"
    
    # Initialize Session State for Drawing History
    if "drawing_chat_history_data" not in st.session_state:
        st.session_state.drawing_chat_history_data = load_history(DRAWING_HISTORY_FILE)
    
    if "current_drawing_session_id" not in st.session_state:
        new_id = str(uuid.uuid4())
        st.session_state.current_drawing_session_id = new_id
        st.session_state.drawing_chat_history_data[new_id] = {
            "title": "새로운 대화",
            "timestamp": datetime.now().isoformat(),
            "messages": []
        }
        st.session_state.rag_chat_messages = []

    # Layout: Spacer L (25%) | Main Content (50%) | Spacer R (10%) | History Sidebar (15%)
    col_spacer_l, col_main, col_spacer_r, col_history = st.columns([0.25, 0.5, 0.1, 0.15])
    
    # Custom CSS for Sidebar Styling (Same as Home)
    st.markdown("""
    <style>
    /* Target the fourth column (History Sidebar) */
    [data-testid="stHorizontalBlock"] > [data-testid="column"]:nth-of-type(4) {
        background-color: #1E1E1E;
        border-left: 1px solid #333;
        padding: 1rem;
        border-radius: 10px;
    }
    [data-testid="column"]:nth-of-type(4) button {
        text-align: left;
    }
    </style>
    """, unsafe_allow_html=True)

    # --- Right Sidebar (History) ---
    with col_history:
        st.markdown("### 채팅 기록")
        
        if st.button("➕ 새 채팅", key="new_drawing_chat", use_container_width=True):
            new_id = str(uuid.uuid4())
            st.session_state.current_drawing_session_id = new_id
            st.session_state.drawing_chat_history_data[new_id] = {
                "title": "새로운 대화",
                "timestamp": datetime.now().isoformat(),
                "messages": []
            }
            st.session_state.rag_chat_messages = []
            st.rerun()
            
        st.markdown("---")
        
        sorted_sessions = sorted(
            st.session_state.drawing_chat_history_data.items(),
            key=lambda x: x[1].get("timestamp", ""),
            reverse=True
        )
        
        for session_id, session_data in sorted_sessions:
            title = session_data.get("title", "대화")
            if session_id == st.session_state.current_drawing_session_id:
                if st.button(f"📂 {title}", key=f"drawing_hist_{session_id}", use_container_width=True, type="primary"):
                    pass
            else:
                if st.button(f"📄 {title}", key=f"drawing_hist_{session_id}", use_container_width=True):
                    st.session_state.current_drawing_session_id = session_id
                    st.session_state.rag_chat_messages = session_data.get("messages", [])
                    st.rerun()
        
        if st.button("🗑️ 기록 삭제", key="del_drawing_hist", use_container_width=True):
            st.session_state.drawing_chat_history_data = {}
            save_history(DRAWING_HISTORY_FILE, {})
            new_id = str(uuid.uuid4())
            st.session_state.current_drawing_session_id = new_id
            st.session_state.drawing_chat_history_data[new_id] = {
                "title": "새로운 대화",
                "timestamp": datetime.now().isoformat(),
                "messages": []
            }
            st.session_state.rag_chat_messages = []
            st.rerun()

    with col_main:
        st.title("도면/스펙 비교")
    
        tab1, tab2 = st.tabs(["📤 문서 업로드", "💬 AI분석"])
    
        with tab1:
        
            if "drawing_uploader_key" not in st.session_state:
                st.session_state.drawing_uploader_key = 0
            
            # High Resolution OCR Toggle
            use_high_res = st.toggle("고해상도 OCR 적용 (도면 미세 글자 추출용)", value=False, help="복잡한 도면의 작은 글씨를 더 정확하게 읽습니다. 분석 시간이 더 오래 걸릴 수 있습니다.")
        

        
            # --- RESUME UI SECTION ---
            # Check for interrupted sessions
            import glob
            resumable_files = []
            if os.path.exists(FILES_DIR):
                for filepath in glob.glob(os.path.join(FILES_DIR, "*")):
                    filename = os.path.basename(filepath)
                    # Check if progress exists
                    if load_progress(filename):
                        resumable_files.append(filename)
            
            files_to_process = []
            
            if resumable_files:
                st.warning(f"⚠️ 중단된 분석 작업이 {len(resumable_files)}건 발견되었습니다.")
                for r_file in resumable_files:
                    progress = load_progress(r_file)
                    processed = progress.get('processed_pages', 0)
                    total = progress.get('total_pages', '?')
                    last_updated = progress.get('last_updated', 'Unknown')
                    
                    # Format timestamp
                    try:
                        dt = datetime.fromisoformat(last_updated)
                        time_str = dt.strftime("%Y-%m-%d %H:%M:%S")
                    except:
                        time_str = last_updated
                    
                    with st.expander(f"📄 {r_file} ({processed}/{total} 페이지 완료) - {time_str}", expanded=True):
                        col_res1, col_res2 = st.columns(2)
                        if col_res1.button(f"▶️ 이어서 분석하기 (Resume)", key=f"resume_{r_file}"):
                            # Create mock file object
                            local_path = get_temp_file_path(r_file)
                            if os.path.exists(local_path):
                                mock_file = LocalFile(local_path, r_file)
                                files_to_process.append(mock_file)
                                st.session_state.is_resuming = True
                            else:
                                st.error("원본 임시 파일을 찾을 수 없습니다.")
                        
                        if col_res2.button(f"🗑️ 취소 및 삭제 (Discard)", key=f"discard_{r_file}"):
                            delete_progress(r_file)
                            delete_temp_file(r_file)
                            st.rerun()

            uploaded_files = st.file_uploader("PDF 도면, 스펙, 사양서 등을 업로드하세요", accept_multiple_files=True, type=['pdf', 'png', 'jpg', 'jpeg', 'tiff', 'bmp'], key=f"drawing_{st.session_state.drawing_uploader_key}")
        
            if uploaded_files or files_to_process:
                if "analysis_status" not in st.session_state:
                    st.session_state.analysis_status = {}
                
                # If resuming, we skip the "Start" button check or auto-click it
                start_analysis = False
                if files_to_process:
                    start_analysis = True
                    target_files = files_to_process
                elif uploaded_files:
                    # Immediate DRM Check
                    drm_files = [f.name for f in uploaded_files if is_drm_protected(f)]
                    if drm_files:
                        st.session_state.drm_error_message = f"⛔ 다음 파일들은 DRM으로 보호되어 있어 업로드할 수 없습니다: {', '.join(drm_files)}. 파일 목록이 초기화되었습니다."
                        st.session_state.drawing_uploader_key += 1
                        st.rerun()

                    if st.button("업로드 및 분석 시작"):
                        start_analysis = True
                        target_files = uploaded_files
                
                if start_analysis:
                    blob_service_client = get_blob_service_client()
                    container_client = blob_service_client.get_container_client(CONTAINER_NAME)
                    doc_intel_manager = get_doc_intel_manager()
                    search_manager = get_search_manager()
                    
                    progress_bar = st.progress(0)
                    status_text = st.empty()
                    
                    total_files = len(target_files)
                    
                    for idx, file in enumerate(target_files):
                        try:
                            # Normalize filename to NFC (to match search query logic)
                            import unicodedata
                            safe_filename = unicodedata.normalize('NFC', file.name)
                            
                            # Save to temp dir for resume capability (only if it's a fresh upload)
                            if not isinstance(file, LocalFile):
                                save_uploaded_file_temp(file, safe_filename)
                            
                            # Initialize status
                            st.session_state.analysis_status[safe_filename] = {
                                "status": "Extracting",
                                "total_pages": 0,
                                "processed_pages": 0,
                                "chunks": {},
                                "error": None
                            }
                            
                            status_text.text(f"처리 중 ({idx+1}/{total_files}): {safe_filename}")
                            
                            blob_path = f"{user_folder}/drawings/{safe_filename}"
                            # 2. Upload to Azure Blob Storage
                            status_text.text(f"업로드 중 ({idx+1}/{total_files}): {file.name}...")
                            blob_client = blob_service_client.get_blob_client(container=CONTAINER_NAME, blob=blob_path)
                            
                            # CRITICAL: Reset file pointer to ensure full upload
                            file.seek(0)
                            blob_client.upload_blob(file, overwrite=True)
                            
                            # Verify upload size
                            props = blob_client.get_blob_properties()
                            if props.size != file.size:
                                st.error(f"⚠️ 파일 업로드 크기 불일치! (원본: {file.size}, 업로드됨: {props.size})")
                            else:
                                print(f"DEBUG: Upload verified. Size: {props.size} bytes")

                            # Generate SAS Token for Document Intelligence access
                            sas_token = generate_blob_sas(
                                account_name=blob_service_client.account_name,
                                container_name=CONTAINER_NAME,
                                blob_name=blob_path,
                                account_key=blob_service_client.credential.account_key,
                                permission=BlobSasPermissions(read=True),
                                expiry=datetime.utcnow() + timedelta(hours=1)
                            )
                            blob_url = f"https://{blob_service_client.account_name}.blob.core.windows.net/{CONTAINER_NAME}/{urllib.parse.quote(blob_path)}?{sas_token}"
                            
                            # 3. Analyze with Document Intelligence (Chunked)
                            file.seek(0)
                            pdf_data = file.read()
                            doc = fitz.open(stream=pdf_data, filetype="pdf")
                            total_pages = doc.page_count
                            file.seek(0)
                            
                            status_text.text(f"분석 준비 중 ({idx+1}/{total_files}): {file.name} (총 {total_pages} 페이지)")
                            
                            st.session_state.analysis_status[safe_filename]["total_pages"] = total_pages
                            
                            chunk_size = 50
                            page_chunks = []
                            
                            # RESUME LOGIC: Check for existing progress
                            existing_progress = load_progress(safe_filename)
                            processed_ranges = set()
                            
                            if existing_progress and existing_progress.get('total_pages') == total_pages:
                                st.info(f"🔄 이전 분석 진행 상황을 발견했습니다. ({len(existing_progress.get('page_chunks', []))} 페이지 완료됨) 이어서 진행합니다.")
                                page_chunks = existing_progress.get('page_chunks', [])
                                
                                # Mark processed ranges based on loaded chunks
                                for chunk in page_chunks:
                                    # We need to reconstruct which ranges are done. 
                                    # Since we don't store ranges explicitly in chunks, we can infer or just skip logic below.
                                    # Better approach: Calculate the range this chunk belongs to and mark it.
                                    p_num = chunk['page_number']
                                    # Calculate start page of the chunk this page belongs to
                                    range_start = ((p_num - 1) // chunk_size) * chunk_size + 1
                                    range_end = min(range_start + chunk_size - 1, total_pages)
                                    processed_ranges.add(f"{range_start}-{range_end}")
                                
                                st.session_state.analysis_status[safe_filename]["processed_pages"] = len(page_chunks)
                            
                            for start_page in range(1, total_pages + 1, chunk_size):
                                end_page = min(start_page + chunk_size - 1, total_pages)
                                page_range = f"{start_page}-{end_page}"
                                
                                # Skip if already processed
                                if page_range in processed_ranges:
                                    st.session_state.analysis_status[safe_filename]["chunks"][page_range] = "Ready"
                                    # status_text.text(f"스킵 중 ({idx+1}/{total_files}): {file.name} - 페이지 {page_range} (이미 완료됨)")
                                    continue
                                
                                st.session_state.analysis_status[safe_filename]["chunks"][page_range] = "Extracting"
                                status_text.text(f"분석 중 ({idx+1}/{total_files}): {file.name} - 페이지 {page_range} 분석 중...")
                                
                                # Retry logic for each chunk
                                max_retries = 3
                                for retry in range(max_retries):
                                    try:
                                        chunks = doc_intel_manager.analyze_document(blob_url, page_range=page_range, high_res=use_high_res)
                                        page_chunks.extend(chunks)
                                        
                                        # Save progress immediately
                                        save_progress(safe_filename, page_chunks, total_pages)
                                        
                                        st.session_state.analysis_status[safe_filename]["chunks"][page_range] = "Ready"
                                        st.session_state.analysis_status[safe_filename]["processed_pages"] += len(chunks)
                                        break
                                    except Exception as e:
                                        if retry == max_retries - 1:
                                            st.session_state.analysis_status[safe_filename]["chunks"][page_range] = "Failed"
                                            st.session_state.analysis_status[safe_filename]["error"] = str(e)
                                            raise e
                                        
                                        # Transient error - show friendly message
                                        wait_time = 5 * (retry + 1)
                                        status_text.text(f"⏳ 일시적 지연으로 재연결 중 ({retry+1}/{max_retries}): {file.name} - 페이지 {page_range} (약 {wait_time}초 대기)...")
                                        time.sleep(wait_time)
                            
                            # 4. Indexing
                            st.session_state.analysis_status[safe_filename]["status"] = "Indexing"
                            
                            if len(page_chunks) == 0:
                                st.warning(f"⚠️ 경고: '{file.name}'에서 페이지를 찾을 수 없습니다.")
                            
                            documents_to_index = []
                            for page_chunk in page_chunks:
                                # Create document object for each page
                                # ID must be unique and URL safe. Include page number in ID.
                                import base64
                                page_id_str = f"{blob_path}_page_{page_chunk['page_number']}"
                                doc_id = base64.urlsafe_b64encode(page_id_str.encode('utf-8')).decode('utf-8')
                                
                                document = {
                                    "id": doc_id,
                                    "content": page_chunk['content'],
                                    "content_exact": page_chunk['content'],
                                    "metadata_storage_name": f"{safe_filename} (p.{page_chunk['page_number']})",
                                    "metadata_storage_path": f"https://{blob_service_client.account_name}.blob.core.windows.net/{CONTAINER_NAME}/{blob_path}#page={page_chunk['page_number']}",
                                    "metadata_storage_last_modified": datetime.utcnow().isoformat() + "Z",
                                    "metadata_storage_size": file.size,
                                    "metadata_storage_content_type": file.type,
                                    "project": "drawings_analysis",  # Tag for filtering
                                    "title": page_chunk.get('도면명(TITLE)', ''),  # Drawing title
                                    "drawing_no": page_chunk.get('도면번호(DWG. NO.)', ''),  # Drawing number
                                    "page_number": page_chunk['page_number'],  # Page number for filtering
                                    "filename": safe_filename  # Filename for search
                                }
                                documents_to_index.append(document)
                            
                            # Batch upload all pages (50 pages at a time to avoid request size limits)
                            indexing_success = True
                            if documents_to_index:
                                batch_size = 50
                                for i in range(0, len(documents_to_index), batch_size):
                                    batch = documents_to_index[i:i + batch_size]
                                    status_text.text(f"인덱싱 중 ({idx+1}/{total_files}): {safe_filename} - 배치 전송 중 ({i//batch_size + 1}/{(len(documents_to_index)-1)//batch_size + 1})")
                                    success, msg = search_manager.upload_documents(batch)
                                    if not success:
                                        st.error(f"인덱싱 실패 ({file.name}, 배치 {i//batch_size + 1}): {msg}")
                                        indexing_success = False
                                        break
                                
                                # 5. Save Analysis JSON to Blob Storage (Dual Retrieval Strategy)
                                # Only save if indexing was successful
                                if indexing_success:
                                    status_text.text(f"분석 결과 저장 중 ({idx+1}/{total_files}): {safe_filename}...")
                                    search_manager.upload_analysis_json(container_client, user_folder, safe_filename, page_chunks)
                                else:
                                    st.warning(f"⚠️ 인덱싱 실패로 인해 '{safe_filename}'의 분석 결과(JSON)를 저장하지 않았습니다.")
                                    # Delete the original file to prevent orphans
                                    try:
                                        st.warning(f"🧹 인덱싱 실패로 인해 원본 파일 '{safe_filename}'을 삭제합니다.")
                                        blob_client.delete_blob()
                                        st.info("원본 파일 삭제 완료.")
                                    except Exception as e:
                                        st.error(f"원본 파일 삭제 실패: {e}")
                            
                            # Success cleanup
                            if indexing_success:
                                delete_progress(safe_filename)
                                delete_temp_file(safe_filename)
                            
                            st.session_state.analysis_status[safe_filename]["status"] = "Ready"
                            progress_bar.progress((idx + 1) / total_files)
                            
                        except Exception as e:
                            st.error(f"오류 발생 ({file.name}): {str(e)}")
                    
                    status_text.text("모든 작업이 완료되었습니다!")
                    st.success("업로드, 분석 및 인덱싱이 완료되었습니다.")
                    
                    # 성공적으로 완료되면 업로더 초기화
                    st.session_state.drawing_uploader_key += 1
                    time.sleep(2)
                    st.rerun()

            # 📊 분석 모니터링 대시보드
            if "analysis_status" in st.session_state and st.session_state.analysis_status:
                st.divider()
                st.markdown("#### 📊 분석 모니터링 대시보드")
                for filename, info in st.session_state.analysis_status.items():
                    status_color = "green" if info['status'] == "Ready" else "orange" if info['status'] != "Failed" else "red"
                    with st.expander(f":{status_color}[{filename}] - {info['status']}", expanded=(info['status'] != "Ready")):
                        col1, col2 = st.columns([3, 1])
                        with col1:
                            st.write(f"**전체 상태:** {info['status']}")
                            progress = info['processed_pages'] / info['total_pages'] if info['total_pages'] > 0 else 0
                            st.progress(progress)
                            st.write(f"**진행도:** {info['processed_pages']} / {info['total_pages']} 페이지 완료")
                    
                        if info['error']:
                            st.error(f"**최근 오류:** {info['error']}")
                    
                        # 세부 청크 상태
                        if info['chunks']:
                            st.markdown("---")
                            st.caption("🧩 **페이지 청크별 상태**")
                            chunk_cols = st.columns(4)
                            for i, (chunk_range, chunk_status) in enumerate(info['chunks'].items()):
                                with chunk_cols[i % 4]:
                                    if chunk_status == "Ready":
                                        st.success(f"✅ {chunk_range}")
                                    elif chunk_status == "Failed":
                                        st.error(f"❌ {chunk_range}")
                                        # 재시도 버튼 (간소화된 구현)
                                        if st.button("🔄", key=f"retry_{filename}_{chunk_range}", help=f"{chunk_range} 재시도"):
                                            st.info("재시도는 '업로드 및 분석 시작'을 다시 눌러주세요 (멱등성 보장)")
                                    else:
                                        st.info(f"⏳ {chunk_range}")

        with tab2:

        
            # Display analyzed documents
            st.markdown("#### 📋 분석된 문서 목록")
            try:
                blob_service_client = get_blob_service_client()
                container_client = blob_service_client.get_container_client(CONTAINER_NAME)
            
                # List files in user's drawings folder + Admin access to root drawings
                blobs = []
                # User folder
                blobs.extend(list(container_client.list_blobs(name_starts_with=f"{user_folder}/drawings/")))
            
                if user_role == 'admin':
                    # Admin root folder
                    blobs.extend(list(container_client.list_blobs(name_starts_with="drawings/")))
            
                # Deduplicate
                unique_blobs = {b.name: b for b in blobs}.values()
            
                blob_list = []
                available_filenames = []
                for blob in unique_blobs:
                    if not blob.name.endswith('/'):  # Skip folder markers
                        filename = blob.name.split('/')[-1]
                        blob_list.append({
                            'name': filename,
                            'full_name': blob.name,
                            'size': blob.size,
                            'modified': blob.last_modified
                        })
                        available_filenames.append(filename)
            
                # Sort by modified date (most recent first)
                blob_list.sort(key=lambda x: x['modified'], reverse=True)
            
                selected_filenames = []
            
                if blob_list:
                    st.info(f"총 {len(blob_list)}개의 문서가 분석되어 있습니다. 분석할 문서를 선택하세요.")
                
                    # Add "Select All" checkbox
                    def toggle_all():
                        new_state = st.session_state.select_all_files
                        # Update state for ALL files in the list, not just existing keys
                        for b in blob_list:
                            st.session_state[f"chk_{b['name']}"] = new_state

                    select_all = st.checkbox("전체 선택", value=False, key="select_all_files", on_change=toggle_all)
                
                    # Display as expandable list
                    with st.expander("📄 문서 목록 및 선택", expanded=True):
                        for idx, blob_info in enumerate(blob_list, 1):
                            # Improved column layout for better zoom stability
                            # col0: checkbox (3%), col1: filename (59%), col2: 3 icons (27%), col3: delete+JSON (11%)
                            col0, col1, col2, col3 = st.columns([0.3, 5.9, 2.7, 1.1])
                            
                            with col0:
                                # Checkbox for selection
                                chk_key = f"chk_{blob_info['name']}"
                                if chk_key not in st.session_state:
                                    st.session_state[chk_key] = False
                                
                                is_selected = st.checkbox(f"select_{idx}", key=chk_key, label_visibility="collapsed")
                                if is_selected:
                                    selected_filenames.append(blob_info['name'])
                        
                            with col1:
                                size_mb = blob_info['size'] / (1024 * 1024)
                                st.markdown(f"**{blob_info['name']}** ({size_mb:.2f} MB)")
                        
                            with col2:
                                # 3 action icons in a row
                                icon_c1, icon_c2, icon_c3 = st.columns(3)
                            
                                with icon_c1:
                                    # Download Button
                                    try:
                                        sas_url = generate_sas_url(
                                            blob_service_client, 
                                            CONTAINER_NAME, 
                                            blob_info['full_name'], 
                                            content_disposition="attachment"
                                        )
                                        st.link_button("📥", sas_url, help="다운로드", use_container_width=True)
                                    except Exception as e:
                                        st.error(f"Err: {e}")

                                with icon_c2:
                                    # 2. Rename Button (Popover)
                                    with st.popover("✏️", use_container_width=True):
                                        new_name_input = st.text_input("새 파일명", value=blob_info['name'], key=f"ren_{blob_info['name']}")
                                        if st.button("이름 변경", key=f"btn_ren_{blob_info['name']}"):
                                            if new_name_input != blob_info['name']:
                                                try:
                                                    with st.spinner("이름 변경 및 인덱스 업데이트 중..."):
                                                        # A. Rename Blob
                                                        old_blob_name = blob_info['full_name']
                                                        folder_path = old_blob_name.rsplit('/', 1)[0]
                                                        new_blob_name = f"{folder_path}/{new_name_input}"
                                                    
                                                        source_blob = container_client.get_blob_client(old_blob_name)
                                                        dest_blob = container_client.get_blob_client(new_blob_name)
                                                    
                                                        # Copy
                                                        source_sas = generate_sas_url(blob_service_client, CONTAINER_NAME, old_blob_name)
                                                        dest_blob.start_copy_from_url(source_sas)
                                                    
                                                        # Wait for copy
                                                        for _ in range(20):
                                                            props = dest_blob.get_blob_properties()
                                                            if props.copy.status == "success":
                                                                break
                                                            time.sleep(0.2)
                                                    
                                                        # B. Update Search Index (Preserve OCR Data)
                                                        search_manager = get_search_manager()
                                                        import unicodedata
                                                        safe_old_filename = unicodedata.normalize('NFC', blob_info['name'])
                                                        safe_new_filename = unicodedata.normalize('NFC', new_name_input)
                                                    
                                                        # Find old docs
                                                        results = search_manager.search_client.search(
                                                            search_text="*",
                                                            filter=f"project eq 'drawings_analysis'",
                                                            select=["id", "content", "content_exact", "metadata_storage_name", "metadata_storage_path", "metadata_storage_size", "metadata_storage_content_type"]
                                                        )
                                                    
                                                        docs_to_upload = []
                                                        ids_to_delete = []
                                                    
                                                        for doc in results:
                                                            # Check if this doc belongs to the file (by name prefix)
                                                            # Name format: "{filename} (p.{page})"
                                                            if doc['metadata_storage_name'].startswith(safe_old_filename):
                                                                # Create new doc
                                                                page_suffix = doc['metadata_storage_name'].split(safe_old_filename)[-1] # e.g. " (p.1)"
                                                            
                                                                # New ID
                                                                import base64
                                                                # Extract page number from suffix or path if possible, or just reconstruct
                                                                # Path format: .../filename#page=N
                                                                try:
                                                                    page_num = doc['metadata_storage_path'].split('#page=')[-1]
                                                                    new_page_id_str = f"{new_blob_name}_page_{page_num}"
                                                                    new_doc_id = base64.urlsafe_b64encode(new_page_id_str.encode('utf-8')).decode('utf-8')
                                                                
                                                                    new_doc = {
                                                                        "id": new_doc_id,
                                                                        "content": doc['content'],
                                                                        "content_exact": doc.get('content_exact', doc['content']),
                                                                        "metadata_storage_name": f"{safe_new_filename}{page_suffix}",
                                                                        "metadata_storage_path": f"https://{blob_service_client.account_name}.blob.core.windows.net/{CONTAINER_NAME}/{new_blob_name}#page={page_num}",
                                                                        "metadata_storage_last_modified": datetime.utcnow().isoformat() + "Z",
                                                                        "metadata_storage_size": doc['metadata_storage_size'],
                                                                        "metadata_storage_content_type": doc['metadata_storage_content_type'],
                                                                        "project": "drawings_analysis"
                                                                    }
                                                                    docs_to_upload.append(new_doc)
                                                                    ids_to_delete.append({"id": doc['id']})
                                                                except:
                                                                    pass
                                                                    
                                                        if docs_to_upload:
                                                            search_manager.upload_documents(docs_to_upload)
                                                        if ids_to_delete:
                                                            search_manager.search_client.delete_documents(documents=ids_to_delete)

                                                        # C. Delete old blob
                                                        source_blob.delete_blob()
                                                    
                                                        st.success("이름 변경 완료!")
                                                        time.sleep(1)
                                                        st.rerun()
                                                    
                                                except Exception as e:
                                                    st.error(f"변경 실패: {e}")

                                with icon_c3:
                                    # 3. Re-analyze Button
                                    if st.button("🔄", key=f"reanalyze_{blob_info['name']}", help="재분석 (인덱스 복구)", use_container_width=True):
                                        try:
                                            with st.spinner("재분석 시작... (파일 다운로드 중)"):
                                                # A. Download Blob to memory
                                                blob_client = container_client.get_blob_client(blob_info['full_name'])
                                                download_stream = blob_client.download_blob()
                                                pdf_data = download_stream.readall()
                                            
                                                # B. Count Pages
                                                import fitz
                                                doc = fitz.open(stream=pdf_data, filetype="pdf")
                                                total_pages = doc.page_count
                                            
                                                # C. Initialize Status
                                                if "analysis_status" not in st.session_state:
                                                    st.session_state.analysis_status = {}
                                            
                                                safe_filename = blob_info['name']
                                                st.session_state.analysis_status[safe_filename] = {
                                                    "status": "Extracting",
                                                    "total_pages": total_pages,
                                                    "processed_pages": 0,
                                                    "chunks": {},
                                                    "error": None
                                                }
                                            
                                                # D. Analyze Chunks
                                                doc_intel_manager = get_doc_intel_manager()
                                                search_manager = get_search_manager()
                                                blob_service_client = get_blob_service_client()
                                            
                                                # Generate SAS for Analysis
                                                sas_token = generate_blob_sas(
                                                    account_name=blob_service_client.account_name,
                                                    container_name=CONTAINER_NAME,
                                                    blob_name=blob_info['full_name'],
                                                    account_key=blob_service_client.credential.account_key,
                                                    permission=BlobSasPermissions(read=True),
                                                    expiry=datetime.utcnow() + timedelta(hours=1)
                                                )
                                                # Use relative path for URL construction if needed, but full_name is usually relative to container if listed from container_client?
                                                # container_client.list_blobs returns name relative to container.
                                                blob_url = f"https://{blob_service_client.account_name}.blob.core.windows.net/{CONTAINER_NAME}/{urllib.parse.quote(blob_info['full_name'])}?{sas_token}"
                                            
                                                chunk_size = 50
                                                page_chunks = []
                                            
                                                progress_bar = st.progress(0)
                                                status_text = st.empty()
                                            
                                                for start_page in range(1, total_pages + 1, chunk_size):
                                                    end_page = min(start_page + chunk_size - 1, total_pages)
                                                    page_range = f"{start_page}-{end_page}"
                                                
                                                    st.session_state.analysis_status[safe_filename]["chunks"][page_range] = "Extracting"
                                                    status_text.text(f"재분석 중: {safe_filename} ({page_range})...")
                                                
                                                    # Retry logic
                                                    max_retries = 3
                                                    for retry in range(max_retries):
                                                        try:
                                                            # Use default high_res=False for re-analysis
                                                            chunks = doc_intel_manager.analyze_document(blob_url, page_range=page_range, high_res=False)
                                                            page_chunks.extend(chunks)
                                                            st.session_state.analysis_status[safe_filename]["chunks"][page_range] = "Ready"
                                                            st.session_state.analysis_status[safe_filename]["processed_pages"] += len(chunks)
                                                            break
                                                        except Exception as e:
                                                            if retry == max_retries - 1:
                                                                st.session_state.analysis_status[safe_filename]["chunks"][page_range] = "Failed"
                                                                st.session_state.analysis_status[safe_filename]["error"] = str(e)
                                                                raise e
                                                            time.sleep(5 * (retry + 1))
                                            
                                                # E. Indexing
                                                st.session_state.analysis_status[safe_filename]["status"] = "Indexing"
                                                status_text.text("인덱싱 중...")
                                            
                                                documents_to_index = []
                                                for page_chunk in page_chunks:
                                                    import base64
                                                    # Use full_name (path in container) for ID generation to match upload logic
                                                    page_id_str = f"{blob_info['full_name']}_page_{page_chunk['page_number']}"
                                                    doc_id = base64.urlsafe_b64encode(page_id_str.encode('utf-8')).decode('utf-8')
                                                
                                                    document = {
                                                        "id": doc_id,
                                                        "content": page_chunk['content'],
                                                        "content_exact": page_chunk['content'],
                                                        "metadata_storage_name": f"{safe_filename} (p.{page_chunk['page_number']})",
                                                        "metadata_storage_path": f"https://{blob_service_client.account_name}.blob.core.windows.net/{CONTAINER_NAME}/{blob_info['full_name']}#page={page_chunk['page_number']}",
                                                        "metadata_storage_last_modified": datetime.utcnow().isoformat() + "Z",
                                                        "metadata_storage_size": blob_info['size'],
                                                        "metadata_storage_content_type": "application/pdf",
                                                        "project": "drawings_analysis",
                                                        "title": page_chunk.get('도면명(TITLE)', ''),  # Drawing title
                                                        "drawing_no": page_chunk.get('도면번호(DWG. NO.)', ''),  # Drawing number
                                                        "page_number": page_chunk['page_number'],  # Page number for filtering
                                                        "filename": safe_filename  # Filename for search
                                                    }
                                                    documents_to_index.append(document)
                                            
                                                if documents_to_index:
                                                    batch_size = 50
                                                    for i in range(0, len(documents_to_index), batch_size):
                                                        batch = documents_to_index[i:i + batch_size]
                                                        success, msg = search_manager.upload_documents(batch)
                                                        if not success:
                                                            st.error(f"❌ 인덱스 업로드 실패 (배치 {i//batch_size + 1}): {msg}")
                                                            raise Exception(f"Index upload failed: {msg}")
                                                
                                                    # Save JSON only if upload succeeded
                                                    search_manager.upload_analysis_json(container_client, user_folder, safe_filename, page_chunks)
                                            
                                                st.session_state.analysis_status[safe_filename]["status"] = "Ready"
                                                st.success("재분석 완료! 이제 검색이 가능합니다.")
                                                time.sleep(1)
                                                st.rerun()

                                        except Exception as e:
                                            st.error(f"재분석 실패: {e}")

                                # 3. JSON (Admin only)
                                if user_role == 'admin':
                                    json_key = f"json_data_{blob_info['name']}"
                                
                                    if json_key not in st.session_state:
                                        if st.button("JSON", key=f"gen_json_{blob_info['name']}"):
                                            with st.spinner("..."):
                                                search_manager = get_search_manager()
                                                # Dual Retrieval Strategy: Try Blob first
                                                docs = search_manager.get_document_json_from_blob(container_client, user_folder, blob_info['name'])
                                            
                                                # Fallback to AI Search if Blob JSON not found (for older files)
                                                if not docs:
                                                    st.info("Blob JSON을 찾을 수 없어 AI Search에서 검색합니다...")
                                                    docs = search_manager.get_document_json(blob_info['name'])
                                                
                                                if docs:
                                                    import json
                                                    json_str = json.dumps(docs, ensure_ascii=False, indent=2)
                                                    st.session_state[json_key] = json_str
                                                    st.rerun()
                                                else:
                                                    st.error(f"No Data found for '{blob_info['name']}'")
                                                    # Try one more time without project filter to see if it exists at all
                                                    safe_name = blob_info['name'].replace("'", "''")
                                                    debug_docs = search_manager.search_client.search(
                                                        search_text="*",
                                                        filter=f"search.ismatch('\"{safe_name}*\"', 'metadata_storage_name')",
                                                        select=["metadata_storage_name", "project"],
                                                        top=5
                                                    )
                                                    debug_list = list(debug_docs)
                                                    if debug_list:
                                                        st.warning(f"Found {len(debug_list)} docs without correct project tag. Example: {debug_list[0].get('metadata_storage_name')} (Project: {debug_list[0].get('project')})")
                                                    else:
                                                        st.error("Document not found in index at all.")
                                    else:
                                        # Show download button
                                        json_data = st.session_state[json_key]
                                        st.download_button(
                                            label="💾",
                                            data=json_data,
                                            file_name=f"{blob_info['name']}.json",
                                            mime="application/json",
                                            key=f"dl_json_{blob_info['name']}"
                                        )

                            with col3:
                                if st.button("🗑️ 삭제", key=f"del_{blob_info['name']}"):
                                    try:
                                        # 1. Delete from Blob Storage (Use full_name)
                                        blob_client = container_client.get_blob_client(blob_info['full_name'])
                                        blob_client.delete_blob()
                                    
                                        # 2. Delete from Search Index
                                        search_manager = get_search_manager()
                                    
                                        # Find docs to delete
                                        import unicodedata
                                        safe_filename = unicodedata.normalize('NFC', blob_info['name'])
                                    
                                        # Clean up index (Find ALL pages)
                                        ids_to_delete = []
                                        while True:
                                            results = search_manager.search_client.search(
                                                search_text="*",
                                                filter=f"project eq 'drawings_analysis'",
                                                select=["id", "metadata_storage_name"],
                                                top=1000
                                            )
                                        
                                            batch_ids = []
                                            for doc in results:
                                                # Use NFC normalization for comparison
                                                doc_name = unicodedata.normalize('NFC', doc['metadata_storage_name'])
                                                if doc_name.startswith(safe_filename):
                                                    batch_ids.append({"id": doc['id']})
                                        
                                            if not batch_ids:
                                                break
                                            
                                            search_manager.search_client.delete_documents(documents=batch_ids)
                                            ids_to_delete.extend(batch_ids)
                                        
                                            # If we found less than 1000, we might be done, but to be safe we continue 
                                            # until a search returns no matches for our file.
                                            # Actually, if we delete them, the next search will return different docs.
                                            # So we continue until no more docs match.
                                            if len(batch_ids) == 0:
                                                break
                                    
                                        # Clear JSON state if exists
                                        json_key = f"json_data_{blob_info['name']}"
                                        if json_key in st.session_state:
                                            del st.session_state[json_key]

                                        st.success(f"{blob_info['name']} 삭제 완료")
                                        st.rerun()
                                    except Exception as e:
                                        st.error(f"삭제 실패: {e}")
                else:
                    st.warning("분석된 문서가 없습니다. '문서 업로드 및 분석' 탭에서 문서를 업로드하세요.")
            except Exception as e:
                st.error(f"문서 목록을 불러오는 중 오류 발생: {e}")
        
            st.divider()
        
            # DEBUG: Show selected files
            if user_role == 'admin':
                # st.write(f"DEBUG: Selected Files ({len(selected_filenames)}): {selected_filenames}")
                pass

            # -----------------------------
            # Advanced Search Options (RAG) - Added to match AI Q&A
            # -----------------------------
            with st.expander("⚙️ 고급 검색 옵션 (RAG 설정)", expanded=False):
                c1, c2 = st.columns(2)
                with c1:
                    rag_use_semantic = st.checkbox("시맨틱 랭커 사용", value=True, key="rag_use_semantic", help="의미 기반 검색을 사용하여 정확도를 높입니다.")
                with c2:
                    rag_search_mode_opt = st.radio("검색 모드", ["all (AND)", "any (OR)"], index=1, horizontal=True, key="rag_search_mode", help="any: 키워드 중 하나라도 포함되면 검색 (추천)")
                    rag_search_mode = "all" if "all" in rag_search_mode_opt else "any"
                
                # Output Format Toggle
                st.write("")
                st.markdown("**답변 형식 (Output Format)**")
                output_format = st.radio(
                    "답변 형식을 선택하세요", 
                    ["Table (표)", "Text (텍스트)"], 
                    index=1, 
                    horizontal=True, 
                    label_visibility="collapsed",
                    key="rag_output_format"
                )
        
            # Chat Interface (Similar to main chat but focused)
            if "rag_chat_messages" not in st.session_state:
                st.session_state.rag_chat_messages = []
            
            for message in st.session_state.rag_chat_messages:
                with st.chat_message(message["role"]):
                    st.markdown(message["content"])
                    if "citations" in message and message["citations"]:
                        st.markdown("---")
                        st.caption("📚 **참조 문서:**")
                        for i, citation in enumerate(message["citations"], 1):
                            filepath = citation.get('filepath', 'Unknown')
                            # Use pre-generated final_url if available, otherwise generate one
                            display_url = citation.get('final_url')
                            if not display_url:
                                try:
                                    blob_service_client = get_blob_service_client()
                                    display_url = generate_sas_url(
                                        blob_service_client, 
                                        CONTAINER_NAME, 
                                        filepath, 
                                        page=citation.get('page')
                                    )
                                except:
                                    display_url = "#"
                        
                            st.markdown(f"{i}. [{filepath}]({display_url})")

            if prompt := st.chat_input("도면이나 스펙에 대해 질문하세요..."):
                st.session_state.rag_chat_messages.append({"role": "user", "content": prompt})
                with st.chat_message("user"):
                    st.markdown(prompt)
            
                with st.chat_message("assistant"):
                    with st.spinner("분석 중..."):
                        try:
                            chat_manager = get_chat_manager()
                        
                            conversation_history = [
                                {"role": msg["role"], "content": msg["content"]}
                                for msg in st.session_state.rag_chat_messages[:-1]
                            ]
                        
                            # Use 'any' search mode for better recall (find documents even with partial keyword match)
                            # This is important because technical drawings may have specific terms
                            # Filter to only search documents from the drawings folder
                            # Pass selected_filenames for specific file filtering
                            # If selected_filenames is empty (user deselected all), we should probably warn or search nothing.
                            # But for now let's pass it. If empty, the chat manager might search nothing or all depending on logic.
                            # Actually, let's default to all if none selected? No, user explicitly deselected.
                            # Let's pass the list as is.
                        
                            # Note: selected_filenames is defined in the outer scope of the tab
                            current_files = selected_filenames
                        
                            # Construct robust filter expression
                            # Include fallback for documents that might have lost their project tag but are in the drawings folder
                            base_filter = "(project eq 'drawings_analysis' or search.ismatch('/drawings/', 'metadata_storage_path'))"
                        
                            # Note: We used to filter by path here, but OData encoding issues caused 0 results.
                            # Now we pass user_folder to chat_manager for Python-side filtering.

                            # Append Output Format Instruction
                            final_prompt = prompt
                            if "Table" in output_format:
                                final_prompt += "\n\n[OUTPUT INSTRUCTION]: Please summarize the comparison in a **Markdown Table**."
                            else:
                                final_prompt += "\n\n[OUTPUT INSTRUCTION]: Please summarize the comparison in **Structured Markdown Text**. Do NOT use a table."

                            response_text, citations, context, final_filter, search_results = chat_manager.get_chat_response(
                                final_prompt, 
                                conversation_history,
                                search_mode=rag_search_mode,
                                use_semantic_ranker=rag_use_semantic,
                                filter_expr=base_filter,
                                available_files=current_files,
                                user_folder=user_folder,
                                is_admin=(user_role == 'admin')
                            )

                            # ---------------------------------------------------------
                            # CRITICAL: Linkify Inline Citations & Escape Tildes
                            # ---------------------------------------------------------
                            
                            # 1. Pre-generate Web Viewer URLs for all citations
                            citation_links = {}
                            processed_citations = [] # Store processed citations with URLs for the bottom list
                            
                            if citations:
                                for cit in citations:
                                    filepath = cit.get('filepath', 'Unknown')
                                    # CRITICAL: Clean filepath from page suffixes like " (p.1)" or " (p.1) (p.1)"
                                    import re
                                    clean_filepath = re.sub(r'\s*\(\s*p\.?\s*\d+\s*\)', '', filepath).strip()
                                    
                                    page = cit.get('page')
                                    url = cit.get('url', '')
                                    
                                    # Generate Web Viewer URL
                                    try:
                                        blob_service_client = get_blob_service_client()
                                        final_url = generate_sas_url(
                                            blob_service_client, 
                                            CONTAINER_NAME, 
                                            clean_filepath, 
                                            page=page
                                        )
                                    except Exception as e:
                                        st.error(f"URL 생성 실패 ({clean_filepath}): {e}")
                                        final_url = "#"
                                    
                                    cit['final_url'] = final_url
                                    processed_citations.append(cit)
                                    
                                    filename = os.path.basename(filepath)
                                    if page:
                                        citation_links[(filename, str(page))] = final_url
                            
                            # 2. Replace text citations with Markdown links (Support both [] and ())
                            if response_text:
                                import re
                                
                                # Reference dictionary to store URLs
                                link_references = {}
                                
                                def replace_citation(match):
                                    content = match.group(1).strip()
                                    
                                    # Remove "문서명:" prefix if present (Common in Korean LLM outputs)
                                    content = re.sub(r'^문서명\s*:\s*', '', content)

                                    # Split by last colon to separate filename and page
                                    if ':' in content:
                                        fname = content.rsplit(':', 1)[0].strip()
                                        p_num = match.group(2)
                                    else:
                                        return match.group(0)
                                        
                                    original_text = match.group(0)
                                    
                                    # Try to find matching citation with fuzzy logic
                                    target_url = None
                                    matched_filename = None
                                    
                                    # 1. Clean LLM filename for comparison
                                    clean_llm = re.sub(r'\.pdf$', '', fname.lower().strip())
                                    
                                    for (k_fname, k_page), url in citation_links.items():
                                        # 2. Clean known filename (remove .pdf and (p.N) suffixes)
                                        clean_known = re.sub(r'\.pdf$', '', k_fname.lower().strip())
                                        clean_known = re.sub(r'\s*\(\s*p\.?\s*\d+\s*\)', '', clean_known).strip()
                                        
                                        # CRITICAL FIX: Skip empty filenames to prevent false positive matches
                                        if not clean_known:
                                            continue

                                        # 3. Match page and filename (fuzzy)
                                        if str(k_page) == str(p_num):
                                            # Exact match or containment (handling "문서명: " residue if regex failed)
                                            if clean_llm == clean_known or clean_llm in clean_known or clean_known in clean_llm:
                                                target_url = url
                                                # CRITICAL FIX: Capture the matched filename to replace text
                                                matched_filename = k_fname
                                                break
                                    
                                    if target_url:
                                        # Use Markdown link for table compatibility
                                        # Escape parentheses in URL to avoid breaking Markdown link
                                        safe_url = target_url.replace('(', '%28').replace(')', '%29')
                                        
                                        # Use Inline Link with Icon and Hover Title for compact display
                                        # Escape parentheses in URL
                                        safe_url = target_url.replace('(', '%28').replace(')', '%29')
                                        
                                        hover_text = original_text
                                        if matched_filename:
                                            # Remove .pdf extension for cleaner display
                                            display_filename = re.sub(r'\.pdf$', '', matched_filename, flags=re.IGNORECASE)
                                            # Use hyphen instead of parentheses to avoid Markdown link issues
                                            hover_text = f"{display_filename} - p.{p_num}"
                                        
                                        # Return Icon Link
                                        return f"[🔗]({safe_url} \"{hover_text}\")"
                                    
                                    # Fallback: Check if original_text contains a URL (e.g. from LLM output)
                                    # Match (http...) or (blob...)
                                    url_match = re.search(r'\((https?://[^)]+|blob:[^)]+)\)', original_text)
                                    if url_match:
                                        fallback_url = url_match.group(1)
                                        safe_url = fallback_url.replace('(', '%28').replace(')', '%29')
                                        
                                        # Construct hover text from captured groups
                                        # fname is from group(1), p_num from group(2) of the main regex
                                        display_filename = re.sub(r'\.pdf$', '', fname.strip(), flags=re.IGNORECASE)
                                        hover_text = f"{display_filename} - p.{p_num}"
                                        
                                        return f"[🔗]({safe_url} \"{hover_text}\")"

                                    return original_text

                                # DEBUG: Show raw response before linkification
                                # st.code(response_text, language="markdown")
                                
                                # Pass 1: Strict match (Double Brackets [[...]]) - NEW STANDARD
                                # Matches [[File: p.1]] or [[File - p.1]]
                                # Also consumes trailing URL if present (e.g. [[...]] (url))
                                pattern_double = r'\[\[(.*?)[:\-]\s*p\.?\s*(\d+)\]\](?:\s*\((?:https?|blob):[^)]+\))?'
                                response_text = re.sub(pattern_double, replace_citation, response_text)

                                # Pass 2: Strict match (Closed parentheses/brackets) - LEGACY FALLBACK
                                # Updated to consume trailing URL if present (e.g. [Title: p.1](url)) to avoid double links
                                # Also updated to allow parentheses in filename (e.g. (File(v1): p.1))
                                # Updated to allow hyphen separator (e.g. [File - p.1])
                                pattern_strict = r'[\[\(]([^\[\]]+?[:\-]\s*p\.?\s*(\d+))[\]\)](?:\s*\((?:https?|blob):[^)]+\))?'
                                response_text = re.sub(pattern_strict, replace_citation, response_text)
                                
                                # Pass 3: Truncated match (Open parentheses at end of string) - Fallback
                                # Only matches if it's at the very end of the string or followed by newline
                                pattern_truncated = r'[\[\(]([^\[\]]+?[:\-]\s*p\.?\s*(\d+))(?:\s*\((?:https?|blob):[^)]+\))?$'
                                response_text = re.sub(pattern_truncated, replace_citation, response_text)

                                # Append Reference Definitions to the end of the response
                                if link_references:
                                    response_text += "\n\n"
                                    for ref_id, url in link_references.items():
                                        response_text += f"[{ref_id}]: {url}\n"

                                # 3. Escape tildes (AFTER linkification to avoid breaking links if they contained tildes, though unlikely in URLs)
                                response_text = response_text.replace('~', '\\~')
                        
                            st.markdown(response_text, unsafe_allow_html=True)
                            
                            # Display Google-like search results (Snippets + Links)
                            if search_results:
                                with st.expander("🔍 검색 결과 및 스니펫 (상위 후보)", expanded=True):
                                    for i, res in enumerate(search_results[:5]): # Show top 5 for clarity
                                        res_name = res.get('metadata_storage_name', 'Unknown')
                                        res_path = res.get('metadata_storage_path', '')
                                    
                                        # Extract snippet from highlights
                                        highlights = res.get('@search.highlights', {})
                                        snippet = highlights.get('content', [""])[0] if highlights else ""
                                        if not snippet:
                                            snippet = res.get('content', '')[:200] + "..."
                                    
                                        # Generate SAS link for the result
                                        try:
                                            # Extract blob path from metadata_storage_path
                                            from urllib.parse import unquote
                                            import re
                                        
                                            if "https://direct_fetch/" in res_path:
                                                # Handle custom direct fetch scheme
                                                path_without_scheme = res_path.replace("https://direct_fetch/", "")
                                                blob_path_part = path_without_scheme.split('#')[0]
                                                blob_path_part = unquote(blob_path_part)
                                            elif CONTAINER_NAME in res_path:
                                                # Handle standard Azure Blob URL
                                                blob_path_part = res_path.split(f"/{CONTAINER_NAME}/")[1].split('#')[0]
                                                blob_path_part = unquote(blob_path_part)
                                            else:
                                                # Fallback or relative path
                                                blob_path_part = res_path
                                        
                                            # CRITICAL FIX: Strip " (p.N)" suffix if present in the path
                                            # This happens if the indexer appended it to the path
                                            blob_path_part = re.sub(r'\s*\(p\.\d+\)$', '', blob_path_part)
                                            
                                            # Generate SAS Token
                                            sas_token = generate_blob_sas(
                                                account_name=blob_service_client.account_name,
                                                container_name=CONTAINER_NAME,
                                                blob_name=blob_path_part,
                                                account_key=blob_service_client.credential.account_key,
                                                permission=BlobSasPermissions(read=True),
                                                expiry=datetime.utcnow() + timedelta(hours=1),
                                                content_disposition="inline",
                                                content_type="application/pdf" # Default to PDF for viewer hint
                                            )
                                            sas_url = f"https://{blob_service_client.account_name}.blob.core.windows.net/{CONTAINER_NAME}/{urllib.parse.quote(blob_path_part, safe='/')}?{sas_token}"

                                            # Use Office Online Viewer for Office files ONLY
                                            # PDF files use direct SAS URL (browser viewer) for better page linking
                                            lower_name = res_name.lower()
                                            if lower_name.endswith(('.pptx', '.ppt', '.docx', '.doc', '.xlsx', '.xls')):
                                                encoded_sas_url = urllib.parse.quote(sas_url)
                                                final_url = f"https://view.officeapps.live.com/op/view.aspx?src={encoded_sas_url}"
                                                link_text = "📄 웹에서 보기 (Office Viewer)"
                                            elif lower_name.endswith('.pdf'):
                                                # Direct SAS URL for PDF (No Google Viewer)
                                                final_url = sas_url
                                                link_text = "📄 문서 열기 (새 탭)"
                                            else:
                                                final_url = sas_url
                                                link_text = "📄 문서 열기 (새 탭)"

                                        except:
                                            final_url = "#"
                                            link_text = "링크 생성 실패"

                                        st.markdown(f"**{i+1}. {res_name}**")
                                        st.write(f"_{snippet}_")
                                        if final_url != "#":
                                            st.markdown(f"[{link_text}]({final_url})")
                                        st.divider()

                            if processed_citations:
                                st.markdown("---")
                                st.caption("📚 **참조 문서 (페이지별 링크):**")
                                
                                # Group citations by filename
                                from collections import defaultdict
                                pages_by_file = defaultdict(set)
                                
                                for cit in processed_citations:
                                    fp = cit.get('filepath', 'Unknown')
                                    pg = cit.get('page')
                                    
                                    # Clean filepath
                                    clean_fp = re.sub(r'\s*\(\s*p\.?\s*\d+\s*\)', '', fp).strip()
                                    
                                    if pg:
                                        try:
                                            pg_int = int(pg)
                                            pages_by_file[clean_fp].add(pg_int)
                                        except:
                                            pass
                                    else:
                                        # Ensure file is listed even if no specific page
                                        if clean_fp not in pages_by_file:
                                            pages_by_file[clean_fp] = set()

                                # Display grouped citations
                                for i, (fp, pages) in enumerate(sorted(pages_by_file.items()), 1):
                                    filename = os.path.basename(fp)
                                    
                                    # Generate Doc URL (Page 1)
                                    try:
                                        blob_service_client = get_blob_service_client()
                                        doc_url = generate_sas_url(
                                            blob_service_client, 
                                            CONTAINER_NAME, 
                                            fp, 
                                            page=1
                                        )
                                    except:
                                        doc_url = "#"
                                    
                                    # Base line: Document Title
                                    line = f"**{i}. [{filename}]({doc_url})**"
                                    
                                    # Append Page Links
                                    sorted_pages = sorted(pages)
                                    if sorted_pages:
                                        page_links = []
                                        for p in sorted_pages:
                                            try:
                                                p_url = generate_sas_url(
                                                    blob_service_client, 
                                                    CONTAINER_NAME, 
                                                    fp, 
                                                    page=p
                                                )
                                                page_links.append(f"[p.{p}]({p_url})")
                                            except:
                                                pass
                                        
                                        if page_links:
                                            line += " — " + " · ".join(page_links)
                                    
                                    st.markdown(line)
                            
                            # Debug: Show Citation Links (Hidden by default)
                            # with st.expander("🔍 링크 디버깅 (Debug Links)", expanded=False):
                            #     st.write("Citation Links Keys:", list(citation_links.keys()))
                            #     st.write("Processed Citations:", processed_citations)
                        
                            # Debug: Show Context
                            with st.expander("🔍 검색된 컨텍스트 확인 (Debug Context)", expanded=False):
                                if final_filter:
                                    st.caption(f"**OData Filter:** `{final_filter}`")
                                if search_results:
                                    st.caption(f"**Search Results:** {len(search_results)} chunks found")
                                st.text_area("LLM에게 전달된 원문 데이터", value=context, height=300)

                            st.session_state.rag_chat_messages.append({
                                "role": "assistant",
                                "content": response_text,
                                "citations": citations,
                                "context": context
                            })
                            
                            # --- Auto-Save History ---
                            current_id = st.session_state.current_drawing_session_id
                            current_title = st.session_state.drawing_chat_history_data[current_id]["title"]
                            if current_title == "새로운 대화" and len(st.session_state.rag_chat_messages) > 0:
                                new_title = get_session_title(st.session_state.rag_chat_messages)
                                st.session_state.drawing_chat_history_data[current_id]["title"] = new_title
                            
                            st.session_state.drawing_chat_history_data[current_id]["messages"] = st.session_state.rag_chat_messages
                            st.session_state.drawing_chat_history_data[current_id]["timestamp"] = datetime.now().isoformat()
                            save_history(DRAWING_HISTORY_FILE, st.session_state.drawing_chat_history_data)
                            
                            st.rerun()


                        except Exception as e:
                            st.error(f"오류: {e}")
                            import traceback
                            st.code(traceback.format_exc())
elif menu == "디버그 (Debug)":
    st.title("🕵️‍♂️ RAG Deep Diagnostic Tool (Integrated)")
    
    # Check if admin
    if user_role != 'admin':
        st.error("Admin access required.")
        st.stop()

    search_manager = get_search_manager()
    blob_service_client = get_blob_service_client()
    container_client = blob_service_client.get_container_client(CONTAINER_NAME)

    # Fetch list of files for selection (Filter for drawings only)
    blob_list = []
    try:
        blobs = container_client.list_blobs()
        for b in blobs:
            # Filter: Must be a file (not folder) AND must be in a 'drawings' folder
            if not b.name.endswith('/') and '/drawings/' in b.name:
                blob_list.append(b.name)
    except Exception as e:
        st.error(f"Failed to list blobs: {e}")
    
    blob_list.sort(key=lambda x: x.split('/')[-1]) # Sort by filename
    
    target_blob = st.selectbox("Select Target File", blob_list)
    
    # Extract filename for search
    if target_blob:
        filename = target_blob.split('/')[-1]
        st.caption(f"Selected Filename for Search: `{filename}`")
    else:
        filename = st.text_input("Target Filename", value="제4권 도면(청주).pdf")

    if st.button("Run Diagnostics"):
        st.divider()
        
        # 1. Index Inspection
        st.subheader("1. Index Inspection")
        
        # Search for ALL pages (including chunks like "filename (p.1)")
        try:
            # Search Strategy: Look for documents where metadata_storage_name starts with the filename
            # This will catch both the main file and all page chunks
            import unicodedata
            norm_filename = unicodedata.normalize('NFC', filename)
            
            # Use text search for the filename and then filter client-side
            results = search_manager.search_client.search(
                search_text=f"\"{filename}\"",
                search_mode="all",
                select=["id", "metadata_storage_name", "metadata_storage_path", "project", "content"],
                top=1000  # Increase to capture all pages
            )
            
            # Filter to get documents that start with our filename (including page chunks)
            results = [
                doc for doc in results 
                if unicodedata.normalize('NFC', doc.get('metadata_storage_name', '')).startswith(norm_filename)
            ]
            
        except Exception as e:
            st.warning(f"Search failed ({str(e)}). This might indicate an indexing issue.")
            results = []
        
        st.write(f"Found **{len(results)}** documents in index.")
        
        # Show breakdown by type
        if results:
            main_docs = [d for d in results if '(p.' not in d.get('metadata_storage_name', '')]
            page_docs = [d for d in results if '(p.' in d.get('metadata_storage_name', '')]
            st.caption(f"📄 Main file: {len(main_docs)} | 📑 Page chunks: {len(page_docs)}")
        
        if results:
            # Analyze First Result
            first = results[0]
            st.json({
                "First Doc ID": first['id'],
                "Name": first['metadata_storage_name'],
                "Path": first['metadata_storage_path'],
                "Project": first['project']
            })

            # 2. Blob Verification
            st.subheader("2. Blob Verification")
            path = first['metadata_storage_path']
            blob_path = None
            
            if "https://direct_fetch/" in path:
                st.warning("⚠️ Using 'direct_fetch' scheme. This is a virtual path.")
                blob_path = path.replace("https://direct_fetch/", "").split('#')[0]
            elif CONTAINER_NAME in path:
                try:
                    blob_path = path.split(f"/{CONTAINER_NAME}/")[1].split('#')[0]
                    blob_path = urllib.parse.unquote(blob_path)
                except:
                    pass
            
            if blob_path:
                st.write(f"**Extracted Blob Path:** `{blob_path}`")
                blob_client = container_client.get_blob_client(blob_path)
                if blob_client.exists():
                    st.success("✅ Blob exists in storage.")
                else:
                    st.error("❌ Blob DOES NOT exist at this path!")
                    
                    # Search for it
                    st.write("Searching for file in container...")
                    found_blobs = list(container_client.list_blobs(name_starts_with=os.path.dirname(blob_path)))
                    if found_blobs:
                        st.write("Found similar blobs:")
                        for b in found_blobs:
                            st.code(b.name)
                    else:
                        st.warning("No similar blobs found.")
            else:
                st.error("Could not extract blob path from metadata.")

            # 3. List Page Check
            st.subheader("3. List Page Check")
            list_keywords = ["PIPING AND INSTRUMENT DIAGRAM FOR LIST", "DRAWING LIST", "도면 목록"]
            found_list = False
            
            for doc in results:
                # Handle None content safely
                content = doc.get('content')
                if content is None:
                    content = ""
                    st.warning(f"⚠️ Document '{doc['metadata_storage_name']}' has NO CONTENT (NULL).")
                
                content_upper = content.upper()
                if any(k in content_upper for k in list_keywords):
                    st.success(f"✅ Found List Page! Name: `{doc['metadata_storage_name']}`")
                    st.text_area("Content Preview", content[:500], height=150)
                    found_list = True
                    break
            
            if not found_list:
                st.error("❌ List Page NOT found in the top 50 results.")
                st.write("Top 5 Results Content Snippets:")
                for i, doc in enumerate(results[:5]):
                    content_preview = (doc.get('content') or "")[:100]
                    st.text(f"{i+1}. {doc['metadata_storage_name']}: {content_preview}...")

            # 4. Cleanup Tool
            st.divider()
            st.subheader("4. Index Cleanup")
            st.warning("If this document is corrupt (No Content / No Project), you can delete it here.")
            
            if st.button(f"🗑️ Delete ALL {len(results)} found documents from Index"):
                try:
                    # Collect IDs
                    ids_to_delete = [{"id": doc['id']} for doc in results]
                    search_manager.search_client.delete_documents(documents=ids_to_delete)
                    st.success(f"Successfully deleted {len(results)} documents.")
                    st.rerun()
                except Exception as e:
                    st.error(f"Delete failed: {e}")

        else:
            st.error("No documents found in index matching this filename.")
            
            # Debug: List what IS in the index (drawings only)
            st.divider()
            st.subheader("🕵️ Index Content Peek (Top 20 Drawings)")
            try:
                # Get top 20 docs from drawings_analysis project
                peek_results = search_manager.search_client.search(
                    search_text="*",
                    filter="project eq 'drawings_analysis'",
                    select=["metadata_storage_name", "project", "metadata_storage_last_modified"],
                    top=20
                )
                peek_list = list(peek_results)
                if peek_list:
                    st.write(f"Index contains at least {len(peek_list)} documents. Here are the top 20:")
                    peek_data = []
                    for d in peek_list:
                        peek_data.append({
                            "Name": d.get('metadata_storage_name'),
                            "Project": d.get('project'),
                            "Modified": d.get('metadata_storage_last_modified')
                        })
                    st.table(peek_data)
                else:
                    st.error("⚠️ The Index appears to be COMPLETELY EMPTY.")
            except Exception as e:
                st.error(f"Failed to peek index: {e}")

    # -----------------------------
    # 디버깅 도구 (Debug Tools)
    # -----------------------------
    if user_role == 'admin':
        with st.expander("🛠️ 인덱스 및 검색 진단 (Debug Tools)", expanded=False):
            st.warning("이 도구는 검색 문제를 진단하기 위한 것입니다.")
            
            # Secret Inspector
            st.write("### 🔐 자격 증명 확인 (Secret Inspector)")
            def mask_secret(s):
                if not s: return "Not Set"
                if len(s) <= 8: return "*" * len(s)
                return s[:4] + "*" * (len(s)-8) + s[-4:]
            
            secrets_to_check = {
                "AZURE_STORAGE_CONNECTION_STRING": STORAGE_CONN_STR,
                "AZURE_BLOB_CONTAINER_NAME": CONTAINER_NAME,
                "AZURE_OPENAI_ENDPOINT": AZURE_OPENAI_ENDPOINT,
                "AZURE_OPENAI_KEY": AZURE_OPENAI_KEY,
                "AZURE_SEARCH_ENDPOINT": SEARCH_ENDPOINT,
                "AZURE_SEARCH_KEY": SEARCH_KEY,
                "AZURE_TRANSLATOR_KEY": TRANSLATOR_KEY,
                "AZURE_DOC_INTEL_ENDPOINT": AZURE_DOC_INTEL_ENDPOINT,
                "AZURE_DOC_INTEL_KEY": AZURE_DOC_INTEL_KEY
            }
            
            import pandas as pd
            secret_data = []
            for k, v in secrets_to_check.items():
                secret_data.append({"Secret Key": k, "Status": "✅ Loaded" if v else "❌ Missing", "Value (Masked)": mask_secret(v)})
            
            st.table(pd.DataFrame(secret_data))
            
            st.write("---")
            
            if st.button("🔍 인덱스 상태 및 검색 테스트 실행"):
                try:
                    search_manager = get_search_manager()
                    client = search_manager.search_client
                    
                    st.write("### 1. 인덱스 문서 확인 (project='drawings_analysis')")
                    results = client.search(search_text="*", filter="project eq 'drawings_analysis'", select=["id", "metadata_storage_name", "project"], top=20)
                    
                    docs = list(results)
                    st.write(f"Found {len(docs)} docs with project='drawings_analysis'")
                    
                    if docs:
                        for doc in docs:
                            st.code(f"ID: {doc['id']}\nName: {doc['metadata_storage_name']}\nProject: {doc['project']}")
                    
                    st.write("---")
                    st.write("### 1-B. 인덱스 문서 확인 (전체 - 필터 없음)")
                    results_all = client.search(search_text="*", select=["id", "metadata_storage_name", "project"], top=20)
                    docs_all = list(results_all)
                    st.write(f"Found {len(docs_all)} docs in total (top 20)")
                    for doc in docs_all:
                        proj = doc.get('project', 'None')
                        st.code(f"Name: {doc['metadata_storage_name']}\nProject: {proj}")
                    
                    st.write("---")
                    st.write("### 2. 키워드 검색 테스트 ('foundation loading data')")
                    search_results = client.search(search_text="foundation loading data", filter="project eq 'drawings_analysis'", top=5, select=["metadata_storage_name", "content"])
                    search_docs = list(search_results)
                    
                    st.write(f"검색 결과: {len(search_docs)}개")
                    for doc in search_docs:
                        st.text(f"Match: {doc['metadata_storage_name']}")
                        st.caption(f"Content: {doc['content'][:200]}...")
                    
                    st.write("---")
                    st.write("### 3. 와일드카드 검색 테스트 ('*')")
                    wild_results = client.search(search_text="*", filter="project eq 'drawings_analysis'", top=5, select=["metadata_storage_name", "content"])
                    wild_docs = list(wild_results)
                    
                    st.write(f"검색 결과: {len(wild_docs)}개")
                    for doc in wild_docs:
                        st.text(f"Match: {doc['metadata_storage_name']}")
                        st.caption(f"Content: {doc['content'][:200]}...")
                        
                except Exception as e:
                    st.error(f"진단 중 오류 발생: {str(e)}")
                    st.code(str(e))
            
            st.write("---")
            st.write("### 🔍 인덱스 데이터 확인")
            if st.button("📑 인덱스된 모든 파일명 보기"):
                with st.spinner("인덱스 조회 중..."):
                    try:
                        search_manager = get_search_manager()
                        # Get all docs (limit to top 1000 to be safe)
                        results = search_manager.search("*", select=["metadata_storage_name"], top=1000)
                        indexed_files = set()
                        for res in results:
                            # Remove page suffix (p.N) to get base filename
                            name = res['metadata_storage_name']
                            base_name = name.split(' (p.')[0]
                            indexed_files.add(base_name)
                        
                        st.write(f"총 {len(indexed_files)}개의 파일이 인덱스에서 발견되었습니다.")
                        st.dataframe(list(indexed_files), use_container_width=True)
                    except Exception as e:
                        st.error(f"조회 실패: {e}")

            st.write("---")
            st.write("### 🧪 사용자 지정 검색 테스트")
            debug_query = st.text_input("검색어 입력 (예: filter element)", key="debug_query")
            if st.button("검색 테스트 실행", key="run_debug_search"):
                if debug_query:
                    try:
                        search_manager = get_search_manager()
                        client = search_manager.search_client
                        
                        st.write(f"Query: '{debug_query}'")
                        # Use 'any' search mode to match behavior
                        results = client.search(
                            search_text=debug_query, 
                            filter="project eq 'drawings_analysis'", 
                            search_mode="any",
                            select=["metadata_storage_name", "content"],
                            top=10
                        )
                        docs = list(results)
                        st.write(f"검색 결과: {len(docs)}개")
                        
                        if docs:
                            for doc in docs:
                                st.text(f"Match: {doc['metadata_storage_name']}")
                                st.caption(f"Content: {doc['content'][:200]}...")
                        else:
                            st.warning("검색 결과가 없습니다.")
                    except Exception as e:
                        st.error(f"검색 오류: {e}")

            st.write("---")
            st.write("### ⚠️ 인덱스 초기화")
            if st.button("🗑️ 모든 도면 데이터 삭제 (Index & Blob)", type="primary"):
                try:
                    # 1. Delete all blobs in any drawings/, json/ folder (Global reset)
                    blob_service_client = get_blob_service_client()
                    container_client = blob_service_client.get_container_client(CONTAINER_NAME)
                    
                    # List all blobs and filter for drawings or json
                    blobs = container_client.list_blobs()
                    deleted_blobs = 0
                    for blob in blobs:
                        if '/drawings/' in blob.name or blob.name.startswith('drawings/') or '/json/' in blob.name or blob.name.startswith('json/'):
                            container_client.delete_blob(blob.name)
                            deleted_blobs += 1
                    
                    # 2. Delete all docs in index with project='drawings_analysis'
                    search_manager = get_search_manager()
                    
                    deleted_total = 0
                    while True:
                        results = search_manager.search_client.search(
                            search_text="*",
                            filter="project eq 'drawings_analysis'",
                            select=["id"],
                            top=1000
                        )
                        ids_to_delete = [{"id": doc['id']} for doc in results]
                        if not ids_to_delete:
                            break
                            
                        search_manager.search_client.delete_documents(documents=ids_to_delete)
                        deleted_total += len(ids_to_delete)
                        if len(ids_to_delete) < 1000:
                            break
                    
                    st.success(f"모든 도면 데이터가 삭제되었습니다. (Blob 삭제 완료, Index {deleted_total}개 삭제 완료) 이제 파일을 다시 업로드하세요.")
                    st.rerun()
                except Exception as e:
                    st.error(f"초기화 실패: {e}")

            if st.button("🧹 '페이지 번호 없는' 중복 데이터 정리 (권장)", help="인덱스에서 (p.N) 형식이 아닌 잘못된 데이터를 찾아 삭제합니다."):
                try:
                    search_manager = get_search_manager()
                    results = search_manager.search_client.search(
                        search_text="*",
                        filter="project eq 'drawings_analysis'",
                        select=["id", "metadata_storage_name"],
                        top=1000
                    )
                    
                    ids_to_delete = []
                    count = 0
                    for doc in results:
                        name = doc['metadata_storage_name']
                        # Delete if it doesn't contain "(p." (standard page suffix)
                        if "(p." not in name:
                            ids_to_delete.append({"id": doc['id']})
                            count += 1
                    
                    if ids_to_delete:
                        search_manager.search_client.delete_documents(documents=ids_to_delete)
                        st.success(f"정리 완료! {count}개의 중복/잘못된 문서를 삭제했습니다.")
                        st.rerun()
                    else:
                        st.info("삭제할 잘못된 데이터가 없습니다. 인덱스가 깨끗합니다! ✨")
                        
                except Exception as e:
                    st.error(f"정리 중 오류 발생: {e}")

            if st.button("🏷️ 누락된 'drawings_analysis' 태그 복구", help="드로잉 폴더에 있지만 프로젝트 태그가 없는 문서를 찾아 태그를 추가합니다."):
                try:
                    search_manager = get_search_manager()
                    # Search for all docs with missing project tag and filter path in Python
                    results = search_manager.search_client.search(
                        search_text="*",
                        filter="(project eq null)",
                        select=["id", "metadata_storage_name", "metadata_storage_path", "content", "content_exact", "metadata_storage_last_modified", "metadata_storage_size", "metadata_storage_content_type"],
                        top=10000 # Increase to cover all docs
                    )
                    
                    docs_to_fix = []
                    for doc in results:
                        # Filter by path in Python
                        if '/drawings/' in doc.get('metadata_storage_path', ''):
                            doc['project'] = 'drawings_analysis'
                            docs_to_fix.append(doc)
                    
                    if docs_to_fix:
                        success, msg = search_manager.upload_documents(docs_to_fix)
                        if success:
                            st.success(f"복구 완료! {len(docs_to_fix)}개의 문서에 'drawings_analysis' 태그를 추가했습니다.")
                            st.rerun()
                        else:
                            st.error(f"복구 실패: {msg}")
                    else:
                        st.info("태그를 복구할 문서가 없습니다.")
                except Exception as e:
                    st.error(f"태그 복구 중 오류 발생: {e}")

            if st.button("📊 인덱스 통계 확인", help="프로젝트별 문서 개수를 확인합니다."):
                try:
                    search_manager = get_search_manager()
                    
                    # Count drawings_analysis
                    drawings_res = search_manager.search_client.search(
                        search_text="*",
                        filter="project eq 'drawings_analysis'",
                        include_total_count=True,
                        top=0
                    )
                    drawings_count = drawings_res.get_count()
                    
                    # Count others (likely standard indexed)
                    others_res = search_manager.search_client.search(
                        search_text="*",
                        filter="project eq null",
                        include_total_count=True,
                        top=0
                    )
                    others_count = others_res.get_count()
                    
                    st.write(f"**도면 분석 데이터 (drawings_analysis):** {drawings_count}개")
                    st.write(f"**일반 문서 데이터 (Standard Indexer):** {others_count}개")
                    
                    # Check Standard Indexer Status
                    st.divider()
                    st.write("**표준 인덱서 (Standard Indexer) 상태 확인:**")
                    # Try common indexer names
                    for idx_name in ["pdf-indexer", "indexer-all", "indexer-drawings"]:
                        try:
                            status = search_manager.indexer_client.get_indexer_status(idx_name)
                            last_res = status.last_result
                            if last_res:
                                st.write(f"- `{idx_name}`: {last_res.status} (성공: {last_res.item_count}, 실패: {last_res.failed_item_count})")
                                if last_res.failed_item_count > 0:
                                    with st.expander(f"❌ {idx_name} 에러 상세 보기"):
                                        for err in last_res.errors[:5]:
                                            st.error(f"문서: {err.key}\n에러: {err.message}")
                            else:
                                st.write(f"- `{idx_name}`: 실행 기록 없음")
                        except:
                            pass

                    if drawings_count == 0 and others_count > 0:
                        st.warning("도면 데이터가 하나도 없습니다. 인덱싱 과정에 문제가 있을 수 있습니다.")
                except Exception as e:
                    st.error(f"통계 확인 중 오류 발생: {e}")

            with st.expander("🔍 인덱스 상세 진단 도구", expanded=False):
                st.caption("인덱스에 저장된 실제 파일명과 태그를 직접 확인합니다.")
                
                # Add search input for specific file diagnosis (Outside button for persistence)
                diag_query = st.text_input("진단할 파일명 검색 (일부만 입력 가능)", value="", key="diag_query")
                diag_path_filter = st.checkbox("'/drawings/' 경로만 보기", value=True, key="diag_path_filter")
                
                if st.button("📋 진단 실행 (최근 100개)"):
                    try:
                        search_manager = get_search_manager()
                        
                        # Use a more inclusive search for diagnosis
                        # If query is provided, use it as search_text. If not, use *
                        results = search_manager.search_client.search(
                            search_text=diag_query if diag_query else "*",
                            select=["metadata_storage_name", "project", "metadata_storage_path"],
                            top=1000 # Increase for better diagnosis
                        )
                        
                        dump_data = []
                        for doc in results:
                            name = doc.get('metadata_storage_name', '')
                            path = doc.get('metadata_storage_path', '')
                            
                            if diag_path_filter and '/drawings/' not in path:
                                continue
                                
                            dump_data.append({
                                "Name": name,
                                "Project": doc.get('project'),
                                "Path": path
                            })
                        
                        if dump_data:
                            st.write(f"검색 결과: {len(dump_data)}개의 문서 발견")
                            st.table(dump_data)
                        else:
                            st.warning("검색 결과가 없습니다. 파일명이 인덱스에 존재하지 않거나 필터에 걸러졌을 수 있습니다.")
                            
                        # Extra check: Search by path only if query failed
                        if diag_query and not dump_data:
                            st.info(f"'{diag_query}'로 검색된 결과가 없어 경로 기반으로 다시 찾습니다...")
                            # Use startswith on metadata_storage_path (SimpleField/Filterable)
                            # We don't know the full prefix, but we can try to find anything in drawings
                            path_results = search_manager.search_client.search(
                                search_text="*",
                                filter="startswith(metadata_storage_path, 'https://')", # Broad filter
                                select=["metadata_storage_name", "project", "metadata_storage_path"],
                                top=5000 # Increase to cover more docs
                            )
                            # Filter for '/drawings/' in Python for maximum reliability
                            path_data = [
                                {"Name": d['metadata_storage_name'], "Project": d['project'], "Path": d['metadata_storage_path']} 
                                for d in path_results 
                                if '/drawings/' in d.get('metadata_storage_path', '')
                            ]
                            if path_data:
                                st.write("'/drawings/' 경로에서 발견된 파일들 (최근 100개 중):")
                                st.table(path_data[:20])
                            else:
                                st.error("'/drawings/' 경로에서 문서를 찾을 수 없습니다. 인덱서가 해당 폴더를 스캔하지 않았을 수 있습니다.")
                                
                    except Exception as e:
                        st.error(f"진단 중 오류 발생: {e}")
                        results_list = list(results)
                        st.info(f"검색된 청크(Chunk) 수: {len(results_list)}개")
                        
                        total_chars = 0
                        for i, doc in enumerate(results_list):
                            content = doc.get('content', '')
                            char_count = len(content)
                            total_chars += char_count
                            
                            with st.expander(f"Chunk {i+1}: {doc.get('metadata_storage_name')} ({char_count}자)"):
                                st.code(content[:1000] + ("..." if len(content) > 1000 else ""))
                        
                        st.divider()
                        st.metric("총 글자 수 (Total Characters)", f"{total_chars:,}")
                        est_tokens = int(total_chars / 4)
                        st.metric("예상 토큰 수 (Estimated Tokens)", f"{est_tokens:,}")
                        
                        if est_tokens > 5000:
                            st.warning(f"⚠️ 토큰 수가 많습니다 ({est_tokens} > 5000). AI 답변 생성 시 'Token Limit Exceeded' 오류가 발생할 수 있습니다.")
                        else:
                            st.success(f"✅ 토큰 수가 적절합니다 ({est_tokens}).")
                            
                    except Exception as e:
                        st.error(f"분석 실패: {e}")
            
            if st.button("🗑️ 대화 초기화", key="clear_rag_chat"):
                st.session_state.rag_chat_messages = []
                st.rerun()

    st.markdown("---")

if menu == "엑셀데이터 자동추출":
    # Integrated Excel Tool
    excel_manager.render_excel_tool()

if menu == "사진대지 자동작성":
    st.caption("건설 현장 사진을 업로드하여 Excel 사진대지를 자동으로 생성합니다.")
    
    # Embed Photo Log app via iframe
    st.components.v1.iframe(
        src="https://photo-log-a0215.web.app/",
        height=800,
        scrolling=True
    )

if menu == "작업계획 및 투입비 자동작성":
    st.caption("작업 계획을 수립하고 투입비를 자동으로 산출합니다.")
    
    # Embed Work Schedule app via iframe
    st.components.v1.iframe(
        src="https://workschedule-7b1cf.web.app/",
        height=800,
        scrolling=True
    )

if menu == "관리자 설정":
    # st.subheader("⚙️ 관리자 설정") - Removed to avoid duplication
    st.info("Azure AI Search 리소스를 초기화하거나 상태를 확인합니다.")
    
    # 인덱싱 대상 폴더 설정
    # 폴더 목록 가져오기
    folder_options = ["(전체)"]
    try:
        blob_service_client = get_blob_service_client()
        container_client = blob_service_client.get_container_client(CONTAINER_NAME)
        # walk_blobs를 사용하여 최상위 폴더만 조회
        for blob in container_client.walk_blobs(delimiter='/'):
            if blob.name.endswith('/'):
                folder_options.append(blob.name.strip('/'))
    except Exception as e:
        st.warning(f"폴더 목록을 가져오지 못했습니다: {e}")
        folder_options.append("GULFLNG") # Fallback

    # 기본값 설정 (GULFLNG가 있으면 그걸로, 없으면 전체)
    default_idx = 0
    if "GULFLNG" in folder_options:
        default_idx = folder_options.index("GULFLNG")

    selected_folder = st.selectbox(
        "인덱싱 대상 폴더 선택", 
        folder_options, 
        index=default_idx,
        help="인덱싱할 프로젝트 폴더를 선택하세요."
    )
    
    
    # '(전체)' 선택 시 None으로 처리
    target_folder = None if selected_folder == "(전체)" else selected_folder
    
    st.info("💡 **폴더별 인덱싱**: 각 폴더는 독립적으로 인덱싱됩니다. 다른 폴더의 데이터에 영향을 주지 않습니다.")
    
    # ------------------------------------------------------------------
    # 인덱스 스키마 업데이트 버튼
    # ------------------------------------------------------------------
    st.subheader("📝 인덱스 스키마 관리")
    st.markdown("""
    도면 메타데이터 필드(`title`, `drawing_no`)를 인덱스에 추가합니다.
    
    **주의:** Azure Search는 기존 필드의 타입을 변경할 수 없습니다.  
    스키마 충돌 시 인덱스를 삭제하고 재생성해야 합니다.
    """)
    
    col1, col2 = st.columns(2)
    
    with col1:
        st.markdown("##### 🔄 스키마 업데이트 (기존 데이터 유지)")
        if st.button("스키마 업데이트", help="새 필드만 추가합니다. 기존 필드 변경 시 에러가 발생할 수 있습니다."):
            with st.spinner("인덱스 스키마 업데이트 중..."):
                manager = get_search_manager()
                success, msg = manager.create_index()
                
                if success:
                    st.success(f"✅ {msg}")
                    st.info("""
                    **다음 단계:**
                    1. 기존 파일의 `🔄` 버튼을 클릭해 재분석하거나
                    2. 새 파일을 업로드하세요
                    
                    **새로 추가된 필드:**
                    - `title` (도면명)
                    - `drawing_no` (도면번호)
                    """)
                else:
                    st.error(f"❌ 스키마 업데이트 실패: {msg}")
                    st.warning("⚠️ 필드 충돌이 발생했다면 오른쪽의 '인덱스 삭제 후 재생성'을 사용하세요.")
    
    with col2:
        st.markdown("##### 🗑️ 삭제 후 재생성 (모든 데이터 삭제)")
        st.warning("⚠️ **경고:** 모든 인덱스 데이터가 삭제됩니다!")
        confirm_delete = st.checkbox("삭제 확인 (모든 데이터 삭제)", key="confirm_delete_index")
        if st.button("인덱스 삭제 후 재생성", disabled=not confirm_delete, help="인덱스를 완전히 삭제하고 새로 생성합니다. 파일을 다시 업로드해야 합니다."):
            with st.spinner("인덱스 삭제 및 재생성 중..."):
                manager = get_search_manager()
                
                # 1. 인덱스 삭제
                del_success, del_msg = manager.delete_index()
                if del_success:
                    st.info(f"🗑️ {del_msg}")
                    time.sleep(2)  # 삭제 완료 대기
                    
                    # 2. 인덱스 재생성
                    create_success, create_msg = manager.create_index()
                    if create_success:
                        st.success(f"✅ {create_msg}")
                        st.success("""
                        **인덱스가 성공적으로 재생성되었습니다!**
                        
                        이제 파일을 다시 업로드하거나 재분석하세요.
                        """)
                    else:
                        st.error(f"❌ 재생성 실패: {create_msg}")
                else:
                    st.error(f"❌ 삭제 실패: {del_msg}")
    
    st.divider()
    
    confirm_reset = st.checkbox("위 폴더를 인덱싱 대상으로 설정하고 싶습니다.", key="confirm_reset")
    
    if st.button("🚀 폴더 인덱싱 설정 (Data Source, Indexer)", disabled=not confirm_reset):
        with st.spinner("리소스 생성 중..."):
            manager = get_search_manager()
            
            # 1. Index 확인/생성 (한번만 필요)
            st.write("1. Index 확인 중...")
            success, msg = manager.create_index()
            if success:
                st.success(msg)
            else:
                st.error(msg)
                
            # 2. Data Source (폴더별)
            st.write(f"2. Data Source 생성 중... (폴더: {selected_folder})")
            success, msg, datasource_name = manager.create_data_source(
                SEARCH_DATASOURCE_NAME, 
                STORAGE_CONN_STR, 
                CONTAINER_NAME, 
                query=target_folder,
                folder_name=target_folder
            )
            if success:
                st.success(msg)
            else:
                st.error(msg)
                st.stop()  # Stop execution if datasource creation fails
                
            # 2.5 Skillset (OCR) - Optional
            skillset_name = None
            enable_ocr = st.checkbox("📸 OCR(이미지 텍스트 추출) 활성화", value=False, help="PDF 도면이나 이미지 파일에서 텍스트를 추출합니다. Azure AI Services 키가 필요하며 비용이 발생할 수 있습니다.")
            
            if enable_ocr:
                st.write(f"2.5. Skillset (OCR) 생성 중...")
                # Use Translator Key as Cognitive Services Key (assuming it's a multi-service key)
                cog_key = st.secrets.get("AZURE_TRANSLATOR_KEY", os.environ.get("AZURE_TRANSLATOR_KEY"))
                
                if not cog_key:
                    st.warning("⚠️ Azure AI Services 키(AZURE_TRANSLATOR_KEY)가 설정되지 않아 OCR을 건너뜁니다.")
                else:
                    skillset_name = f"skillset-{target_folder}" if target_folder else "skillset-all"
                    success, msg = manager.create_skillset(skillset_name, cog_key)
                    if success:
                        st.success(msg)
                    else:
                        st.error(f"Skillset 생성 실패: {msg}")
                        skillset_name = None # Fallback to no skillset
                
            # 3. Indexer (폴더별)
            st.write(f"3. Indexer 생성 중... (폴더: {selected_folder})")
            # 기존 인덱서 삭제 (같은 폴더의 이전 설정 제거)
            manager.delete_indexer(target_folder)
            success, msg, indexer_name = manager.create_indexer(target_folder, datasource_name, skillset_name=skillset_name)
            if success:
                st.success(msg)
                st.info(f"✅ '{selected_folder}' 폴더에 대한 인덱싱 설정이 완료되었습니다. 아래 '인덱서 수동 실행'을 눌러 인덱싱을 시작하세요.")
            else:
                st.error(msg)
    
    st.divider()
    
    # ------------------------------------------------------------------
    # 4. 인덱스 내용 조회 (디버깅용)
    # ------------------------------------------------------------------
    st.subheader("🔍 인덱스 내용 조회 (OCR 확인용)")
    with st.expander("특정 파일의 인덱싱된 내용 확인하기"):
        target_filename = st.text_input("확인할 파일명 (예: drawing.pdf)", help="정확한 파일명을 입력하세요.")
        if st.button("내용 조회"):
            if target_filename:
                manager = get_search_manager()
                with st.spinner("조회 중..."):
                    content = manager.get_document_content(target_filename)
                    st.text_area("인덱싱된 내용 (앞부분 2000자)", content[:2000], height=300)
            else:
                st.warning("파일명을 입력하세요.")

    st.divider()
    
    # 수동 실행 안내 및 확인
    st.info(f"📂 **현재 선택된 폴더**: {selected_folder}")
    st.markdown("수동 인덱서 실행은 선택한 폴더의 새 파일 또는 변경된 파일을 검색 엔진에 반영합니다.")
    
    confirm_run = st.checkbox("위 폴더를 인덱싱하는 것을 확인했으며, 진행하고 싶습니다.", key="confirm_run")
    
    if st.button("▶️ 인덱서 수동 실행", disabled=not confirm_run):
        manager = get_search_manager()
        success, msg = manager.run_indexer(target_folder)
        if success:
            st.success(msg)
            st.info("인덱싱이 시작되었습니다. 아래 '상태 확인' 버튼을 눌러 진행 상황을 모니터링하세요.")
        else:
            st.error(msg)
            
    # Add Delete Indexer Button
    if st.button("🛑 인덱서 삭제 (자동 인덱싱 중지)", help="자동으로 실행되는 인덱서를 삭제하여 중복 인덱싱을 방지합니다."):
        manager = get_search_manager()
        indexer_name = f"indexer-{target_folder}" if target_folder else "indexer-all"
        try:
            manager.indexer_client.delete_indexer(indexer_name)
            st.success(f"인덱서 '{indexer_name}'가 삭제되었습니다. 이제 자동 인덱싱이 중지됩니다.")
        except Exception as e:
            st.error(f"인덱서 삭제 실패: {e}")
            
    st.divider()
    
    col_status, col_refresh = st.columns([3, 1])
    with col_status:
        st.markdown("### 📊 인덱싱 현황 모니터링")
    with col_refresh:
        auto_refresh = st.checkbox("자동 새로고침 (5초)", value=False)

    # 상태 확인 로직 (버튼 클릭 또는 자동 새로고침)
    if st.button("상태 및 진행률 확인") or auto_refresh:
        manager = get_search_manager()
        
        # 1. 소스 파일 개수 확인 (진행률 계산용)
        with st.spinner("소스 파일 개수 계산 중..."):
            total_blobs = manager.get_source_blob_count(STORAGE_CONN_STR, CONTAINER_NAME, folder_path=target_folder)
        
        # 2. 인덱서 상태 확인
        status_info = manager.get_indexer_status(target_folder)
        
        # 상태 언팩
        status = status_info.get("status")
        item_count = status_info.get("item_count", 0)
        failed_count = status_info.get("failed_item_count", 0)
        error_msg = status_info.get("error_message")
        errors = status_info.get("errors", [])
        warnings = status_info.get("warnings", [])
        
        # 3. 인덱스 문서 개수
        doc_count = manager.get_document_count()
        
        # UI 표시
        st.metric(label="총 소스 파일 수", value=f"{total_blobs}개")
        
        # 진행률 계산 (실제 인덱스된 문서 수 기준)
        if total_blobs > 0:
            progress = min(doc_count / total_blobs, 1.0)
        else:
            progress = 0.0
            
        st.progress(progress, text=f"인덱싱 진행률: {int(progress * 100)}% ({doc_count}/{total_blobs})")
        
        # 상태 메시지
        if status == "inProgress":
            st.info(f"⏳ 인덱싱 진행 중... (처리된 문서: {item_count}, 실패: {failed_count})")
            if auto_refresh:
                time.sleep(5)
                st.rerun()
        elif status == "success":
            st.success(f"✅ 인덱싱 완료! (총 인덱스 문서: {doc_count}개)")
        elif status == "error":
            st.error(f"❌ 인덱싱 오류 발생: {error_msg}")
        elif status == "transientFailure":
            st.warning("⚠️ 일시적 오류 발생 (재시도 중...)")
        else:
            st.write(f"상태: {status}")

        # 오류 상세 표시
        if failed_count > 0 or errors:
            st.error(f"❌ 실패한 문서: {failed_count}개")
            with st.expander("🚨 오류 상세 로그 확인", expanded=True):
                for err in errors:
                    st.write(f"- {err}")
        
        if warnings:
            with st.expander("⚠️ 경고 로그 확인"):
                for warn in warnings:
                    st.warning(f"- {warn}")
    
    st.divider()
    
    # ------------------------------------------------------------------
    # 🔍 디버그 툴 - Index Content Peek
    # ------------------------------------------------------------------
    st.subheader("🔍 디버그 툴")
    st.markdown("인덱스에 저장된 문서를 확인하여 스키마 필드가 올바르게 채워졌는지 검증합니다.")
    
    with st.expander("Index Content Peek", expanded=False):
        st.markdown("인덱스에서 최근 문서를 가져와 필드 값을 확인합니다.")
        
        # 프로젝트 필터 옵션
        filter_project = st.text_input(
            "프로젝트 필터 (선택사항)", 
            value="drawings_analysis",
            help="특정 프로젝트의 문서만 조회 (비워두면 모든 문서)"
        )
        
        peek_limit = st.slider("조회할 문서 수", min_value=1, max_value=20, value=5)
        
        if st.button("📄 Peek Index", key="peek_index_btn"):
            with st.spinner("인덱스 조회 중..."):
                try:
                    manager = get_search_manager()
                    
                    # Search with filter
                    if filter_project:
                        results = manager.search_client.search(
                            search_text="*",
                            filter=f"project eq '{filter_project}'",
                            top=peek_limit,
                            select=["id", "content", "title", "drawing_no", "page_number", "filename", "metadata_storage_name", "project"]
                        )
                    else:
                        results = manager.search_client.search(
                            search_text="*",
                            top=peek_limit,
                            select=["id", "content", "title", "drawing_no", "page_number", "filename", "metadata_storage_name", "project"]
                        )
                    
                    docs = list(results)
                    
                    if not docs:
                        st.warning("인덱스에 문서가 없습니다.")
                    else:
                        st.success(f"총 {len(docs)}개 문서를 찾았습니다.")
                        
                        for i, doc in enumerate(docs):
                            with st.expander(f"📄 Document {i+1}: {doc.get('filename', 'N/A')} - Page {doc.get('page_number', 'N/A')}", expanded=(i==0)):
                                # 중요 필드 하이라이트
                                col1, col2 = st.columns(2)
                                with col1:
                                    st.markdown("**핵심 메타데이터:**")
                                    st.json({
                                        "title": doc.get("title"),
                                        "drawing_no": doc.get("drawing_no"),
                                        "page_number": doc.get("page_number"),
                                        "filename": doc.get("filename"),
                                        "project": doc.get("project")
                                    })
                                
                                with col2:
                                    st.markdown("**필드 상태 검증:**")
                                    title_status = "✅" if doc.get("title") else "❌"
                                    drawing_status = "✅" if doc.get("drawing_no") else "❌"
                                    page_status = "✅" if doc.get("page_number") is not None else "❌"
                                    
                                    st.markdown(f"""
                                    - `title`: {title_status} {doc.get("title") or "NULL"}
                                    - `drawing_no`: {drawing_status} {doc.get("drawing_no") or "NULL"}
                                    - `page_number`: {page_status} {doc.get("page_number") if doc.get("page_number") is not None else "NULL"}
                                    """)
                                
                                # Content preview
                                st.markdown("**Content Preview (처음 500자):**")
                                content_preview = doc.get("content", "")[:500]
                                st.text_area("", content_preview, height=150, key=f"content_{i}", disabled=True)
                                
                                # Full JSON
                                with st.expander("전체 JSON 보기"):
                                    st.json(dict(doc))
                
                except Exception as e:
                    st.error(f"인덱스 조회 실패: {e}")
                    import traceback
                    st.code(traceback.format_exc())


if menu == "사용자 설정":
    from modules.user_settings_module import render_user_settings
    render_user_settings(auth_manager)



if menu == "디버그 (Debug)":
    st.title("🔍 Search Debug Tool (Cloud)")
    st.write("Debug Menu Loaded...") # Debug print
    
    # Secrets (Already loaded in app.py as global variables, but we can reuse get_search_manager)
    try:
        search_manager = get_search_manager()
        st.write("Search Manager Loaded.") # Debug print
    except Exception as e:
        st.error(f"Failed to load Search Manager: {e}")
        st.stop()
    



    # ========================================
    st.header("🎯 2단계 검색 테스트 (정확한 키워드 우선)")
    st.info("**목적**: 사용자 입력 그대로 먼저 검색하여 정확한 키워드 매칭을 우선순위로 둡니다.")
    
    test_query = st.text_input(
        "테스트 검색어",
        value="piping and instrument diagram list",
        key="two_stage_query"
    )
    
    test_filename = st.text_input(
        "대상 파일",
        value="제4권 도면(청주).pdf",
        key="two_stage_file"
    )
    
    if st.button("🚀 2단계 검색 실행", type="primary"):
        st.markdown("---")
        
        # Build filter
        filter_expr = None
        if test_filename and test_filename.strip():
            filter_expr = f"search.ismatch('{test_filename}', 'metadata_storage_name')"
        
        # Stage 1: Exact search
        st.subheader("📍 Stage 1: 정확한 키워드 검색 (쿼리 확장 없음)")
        st.code(f"Query: '{test_query}'")
        
        with st.spinner("Stage 1 검색 중..."):
            stage1_results = search_manager.search(
                test_query,  # 원본 그대로
                filter_expr=filter_expr,
                search_mode="any",
                top=50
            )
        
        st.success(f"✅ Stage 1 결과: {len(stage1_results)}개")
        
        if stage1_results:
            st.markdown("**Top 10 결과:**")
            for i, doc in enumerate(stage1_results[:10], 1):
                doc_name = doc.get('metadata_storage_name', 'Unknown')
                content_snippet = doc.get('content', '')[:100].replace('\n', ' ')
                
                # Check if this is page 7
                is_page_7 = "(p.7)" in doc_name
                marker = "🎯 **[TARGET PAGE]** " if is_page_7 else ""
                
                st.markdown(f"{i}. {marker}{doc_name}")
                
                # Detailed view for page 7
                if is_page_7:
                    with st.expander("📄 7페이지 상세 내용"):
                        st.markdown(f"**Content Preview:**")
                        st.text_area("", doc.get('content', '')[:1000], height=200, key=f"p7_content_{i}")
        else:
            st.warning("Stage 1에서 결과를 찾지 못했습니다.")
        
        # Stage 2 simulation
        st.markdown("---")
        st.subheader("📍 Stage 2: 쿼리 확장 검색 (참고용)")
        
        THRESHOLD = 20
        if len(stage1_results) >= THRESHOLD:
            st.info(f"ℹ️ Stage 1에서 {len(stage1_results)}개 결과를 찾았으므로 Stage 2는 **실행되지 않습니다** (threshold: {THRESHOLD})")
        else:
            st.warning(f"⚠️ Stage 1에서 {len(stage1_results)}개만 찾았으므로 Stage 2 쿼리 확장이 필요합니다.")
            
            # Simulate query expansion
            expanded_query = f"{test_query} PIPING INSTRUMENT DIAGRAM LIST INDEX TABLE DRAWING"
            st.code(f"Expanded Query: '{expanded_query}'")
            
            with st.spinner("Stage 2 검색 중..."):
                stage2_results = search_manager.search(
                    expanded_query,
                    filter_expr=filter_expr,
                    search_mode="any",
                    top=50
                )
            
            st.success(f"✅ Stage 2 추가 결과: {len(stage2_results)}개")
    
    st.markdown("---")
    st.markdown("---")
    
    # ========================================
    # 사용자 정의 검색 입력
    # ========================================
    st.header("📝 사용자 지정 검색")
    
    col1, col2 = st.columns([2, 1])
    
    with col1:
        custom_query = st.text_input(
            "검색할 키워드 입력",
            value="piping and instrument diagram list",
            help="검색하고 싶은 키워드를 입력하세요"
        )
    
    with col2:
        custom_top = st.number_input(
            "검색 결과 수",
            min_value=1,
            max_value=200,
            value=50,
            step=10
        )
    
    filename = st.text_input(
        "대상 파일명 (선택사항)",
        value="제4권 도면(청주).pdf",
        help="특정 파일만 검색하려면 입력하세요. 비워두면 전체 인덱스를 검색합니다."
    )
    
    if st.button("🔍 검색 실행", type="primary", use_container_width=True):
        st.markdown("---")
        st.subheader(f"🔎 검색 결과: '{custom_query}'")
        
        with st.spinner("검색 중..."):
            # Build filter
            filter_expr = None
            if filename and filename.strip():
                filter_expr = f"search.ismatch('{filename}', 'metadata_storage_name')"
            
            # Execute search
            results = search_manager.search(
                custom_query,
                filter_expr=filter_expr,
                search_mode="any",
                top=custom_top
            )
            
            st.success(f"✅ **{len(results)}개 결과 발견**")
            
            if len(results) == 0:
                st.warning("검색 결과가 없습니다. 다른 키워드를 시도해보세요.")
            else:
                # Display results
                for i, doc in enumerate(results, 1):
                    doc_name = doc.get('metadata_storage_name', 'Unknown')
                    content = doc.get('content', '')
                    title = doc.get('title', 'No title')
                    
                    with st.expander(f"**{i}. {doc_name}**", expanded=(i <= 3)):
                        st.markdown(f"**Title**: {title}")
                        
                        # Highlight if query keywords are in content
                        content_upper = content.upper()
                        query_upper = custom_query.upper()
                        
                        # Check for keyword presence
                        keywords_found = []
                        for word in query_upper.split():
                            if word in content_upper:
                                keywords_found.append(word)
                        
                        if keywords_found:
                            st.success(f"✅ 키워드 매칭: {', '.join(keywords_found)}")
                        
                        # Content preview
                        st.markdown("**Content Preview (처음 500자):**")
                        st.text_area("", content[:500], height=150, key=f"custom_result_{i}", disabled=True)
                        
                        # Full content
                        with st.expander("전체 내용 보기"):
                            st.text_area("", content, height=400, key=f"custom_full_{i}", disabled=True)

    st.markdown("---")

    st.markdown("---")

    # ---------------------------------------------------------
    # NEW: Keyword Search Debug (Stage 1 Simulation)
    # ---------------------------------------------------------
    st.header("🔍 키워드 검색 및 LLM 컨텍스트 확인 (Keyword Search Debug)")
    st.info("LLM이 특정 정보를 찾지 못할 때, 실제로 검색 엔진이 해당 정보를 찾아내는지 확인하는 도구입니다.")
    
    col_debug_search, col_debug_opts = st.columns([0.7, 0.3])
    
    with col_debug_search:
        debug_keyword = st.text_input("검색할 키워드 입력 (예: 냉각수펌프 전기실)", value="냉각수펌프 전기실")
    
    with col_debug_opts:
        debug_top_k = st.number_input("검색 개수 (Top K)", min_value=1, max_value=50, value=20)
    
    if st.button("🚀 키워드 검색 실행 (Stage 1 Logic)"):
        with st.spinner(f"'{debug_keyword}' 검색 중..."):
            # Use exact same logic as chat_manager Stage 1
            # 1. Sanitize
            import re
            sanitized_query = re.sub(r'\bAND\b', ' ', debug_keyword, flags=re.IGNORECASE)
            sanitized_query = re.sub(r'[&+\-|!(){}\[\]^"~*?:\\]', ' ', sanitized_query)
            sanitized_query = " ".join(sanitized_query.split())
            
            st.write(f"**Sanitized Query:** `{sanitized_query}`")
            
            # 2. Search
            results = search_manager.search(
                sanitized_query,
                use_semantic_ranker=False, # Stage 1 uses standard BM25
                search_mode="all",         # Stage 1 uses AND logic
                top=debug_top_k
            )
            
            if results:
                st.success(f"✅ 총 {len(results)}개의 문서를 찾았습니다.")
                
                debug_data = []
                for rank, res in enumerate(results, 1):
                    name = res.get('metadata_storage_name', 'Unknown')
                    score = res.get('@search.score', 0)
                    content = res.get('content', '')
                    
                    # Highlight keywords in content preview
                    preview = content[:300].replace('\n', ' ') + "..."
                    
                    debug_data.append({
                        "Rank": rank,
                        "Score": f"{score:.4f}",
                        "File": name,
                        "Content Preview": preview
                    })
                
                st.dataframe(pd.DataFrame(debug_data), use_container_width=True)
                
                # Detailed View
                with st.expander("📄 상세 내용 보기 (Top 5)"):
                    for i, res in enumerate(results[:5], 1):
                        st.markdown(f"### {i}. {res.get('metadata_storage_name')}")
                        st.text(res.get('content', '')[:1000])
                        st.markdown("---")
            else:
                st.warning("⚠️ 검색 결과가 없습니다. (No results found)")
                st.markdown("""
                **가능한 원인:**
                1. 문서에 해당 키워드가 정확히 포함되어 있지 않음 (OCR 오류 등)
                2. 'AND' 조건으로 인해 모든 단어가 포함된 문서만 검색됨
                """)
    
    st.markdown("---")

    # ---------------------------------------------------------
    # NEW: Target Page Debug (Why is this page missing?)
    # ---------------------------------------------------------
    st.header("🎯 특정 페이지 검색 누락 원인 분석 (Target Page Debug)")
    st.info("특정 페이지가 검색 결과에 나오지 않을 때, 해당 페이지가 인덱스에 존재하는지, 키워드가 포함되어 있는지 분석합니다.")
    
    col_target_1, col_target_2, col_target_3 = st.columns([2, 2, 1])
    
    with col_target_1:
        target_query = st.text_input("검색 쿼리", value="냉각수펌프 전기실", key="target_debug_query")
    with col_target_2:
        target_filename = st.text_input("파일명 (정확히 입력)", value="제5권 물량내역서(청주).pdf", key="target_debug_file")
    with col_target_3:
        target_page = st.number_input("페이지 번호", value=82, key="target_debug_page")
        
    if st.button("🕵️‍♂️ 페이지 분석 실행"):
        target_doc_name = f"{target_filename} (p.{target_page})"
        st.write(f"**Target Document Name:** `{target_doc_name}`")
        
        # 1. Check if page exists in index
        with st.spinner("인덱스에서 페이지 조회 중..."):
            # Escape single quotes for OData
            safe_doc_name = target_doc_name.replace("'", "''")
            direct_check = search_manager.search(
                "*",
                filter_expr=f"metadata_storage_name eq '{safe_doc_name}'",
                top=1
            )
            
        if not direct_check:
            st.error(f"❌ **페이지가 인덱스에 없습니다!** (`{target_doc_name}`)")
            st.warning("파일명이나 페이지 번호를 확인해주세요. 또는 해당 파일이 인덱싱되지 않았을 수 있습니다.")
        else:
            doc = direct_check[0]
            st.success(f"✅ **페이지가 인덱스에 존재합니다.** (ID: `{doc.get('id', 'N/A')}`)")
            
            raw_content = doc.get('content', '')
            
            # Apply same cleaning logic as Chat Manager
            import re
            cleaned_content = raw_content
            
            # 1. Remove XML comments
            cleaned_content = re.sub(r'<!--.*?-->', '', cleaned_content, flags=re.DOTALL)
            
            # 2. Mark intended line breaks
            LINE_BREAK = "___LB___"
            cleaned_content = re.sub(r'</tr>', LINE_BREAK, cleaned_content, flags=re.IGNORECASE)
            cleaned_content = re.sub(r'<br\s*/?>', LINE_BREAK, cleaned_content, flags=re.IGNORECASE)
            cleaned_content = re.sub(r'</p>', LINE_BREAK, cleaned_content, flags=re.IGNORECASE)
            cleaned_content = re.sub(r'</div>', LINE_BREAK, cleaned_content, flags=re.IGNORECASE)
            
            # 3. Replace cell endings with pipe
            cleaned_content = re.sub(r'</td>', ' | ', cleaned_content, flags=re.IGNORECASE)
            cleaned_content = re.sub(r'</th>', ' | ', cleaned_content, flags=re.IGNORECASE)
            
            # 4. Remove all original newlines
            cleaned_content = cleaned_content.replace('\n', ' ').replace('\r', ' ')
            
            # 5. Remove remaining tags
            cleaned_content = re.sub(r'<[^>]+>', '', cleaned_content)
            
            # 6. Restore intended line breaks
            cleaned_content = cleaned_content.replace(LINE_BREAK, '\n')
            
            # 7. Noise
            cleaned_content = cleaned_content.replace("AutoCAD SHX Text", "").replace("%%C", "Ø")
            
            # 8. Collapse whitespace
            cleaned_content = re.sub(r'[ \t]+', ' ', cleaned_content)
            cleaned_content = re.sub(r'\n\s*\n', '\n\n', cleaned_content)
            cleaned_content = cleaned_content.strip()
            
            st.markdown("### 📄 페이지 내용 (Content Preview)")
            
            tab_clean, tab_raw = st.tabs(["✨ Cleaned (AI가 보는 화면)", "📝 Raw (원본 데이터)"])
            
            with tab_clean:
                st.info("AI에게는 아래와 같이 **표 구조가 정리된 텍스트**가 전달됩니다.")
                st.text_area("Cleaned Content", cleaned_content, height=400)
                
            with tab_raw:
                st.warning("인덱스에 저장된 원본 데이터입니다 (HTML 태그 포함).")
                st.text_area("Raw Content", raw_content, height=400)

            
            # 2. Analyze Keyword Matching (Check against CLEANED content)
            st.markdown("### 🔍 키워드 매칭 분석 (Cleaned Content 기준)")
            keywords = target_query.split()
            
            match_data = []
            content_upper = cleaned_content.upper()
            
            all_matched = True
            for kw in keywords:
                kw_upper = kw.upper()
                count = content_upper.count(kw_upper)
                matched = count > 0
                if not matched:
                    all_matched = False
                
                match_data.append({
                    "Keyword": kw,
                    "Found": "✅ Yes" if matched else "❌ No",
                    "Count": count
                })
            
            st.dataframe(pd.DataFrame(match_data), use_container_width=True)
            
            if all_matched:
                st.success("✅ 모든 키워드가 본문에 포함되어 있습니다. 검색 랭킹 문제일 가능성이 높습니다.")
            else:
                st.error("❌ 일부 키워드가 본문에 없습니다! 이래서 검색이 안 되는 것입니다.")
                st.markdown("""
                **해결 방법:**
                1. **OCR 오류 확인**: 본문 텍스트를 자세히 읽어보세요. 오타가 있나요? (예: `전기실` -> `전 기 실` or `전기슬`)
                2. **동의어 확장**: 사용자가 입력한 단어와 문서에 있는 단어가 다를 수 있습니다.
                """)
                
            # 3. Run actual search to see Rank
            st.markdown("### 📊 실제 검색 랭킹 확인")
            with st.spinner("실제 검색 수행 중..."):
                # Use same logic as Stage 1
                import re
                sanitized_query = re.sub(r'\bAND\b', ' ', target_query, flags=re.IGNORECASE)
                sanitized_query = re.sub(r'[&+\-|!(){}\[\]^"~*?:\\]', ' ', sanitized_query)
                sanitized_query = " ".join(sanitized_query.split())
                
                # Filter by filename to narrow down
                safe_filename = target_filename.replace("'", "''")
                escaped_filename = re.sub(r'([+\-&|!(){}\[\]^"~*?:\\])', r'\\\1', safe_filename)
                file_filter = f"search.ismatch('\"{escaped_filename}\"', 'metadata_storage_name')"
                
                search_results = search_manager.search(
                    sanitized_query,
                    filter_expr=file_filter,
                    use_semantic_ranker=False,
                    search_mode="all",
                    top=200
                )
                
                found_rank = None
                for i, res in enumerate(search_results, 1):
                    if res.get('metadata_storage_name') == target_doc_name:
                        found_rank = i
                        break
                
                if found_rank:
                    st.info(f"ℹ️ 이 페이지는 현재 검색 결과 **{found_rank}위**에 있습니다.")
                else:
                    st.warning("⚠️ 이 페이지는 Top 200 검색 결과에 포함되지 않았습니다.")

    st.markdown("---")

elif menu == "디버그 (Debug)":
    # ========================================
    # 🔍 심층 랭킹 분석 (Deep Ranking Analysis)
    # ========================================
    st.header("🔍 심층 랭킹 분석 (Deep Ranking Analysis)")
    st.info("검색어에 대해 각 페이지가 왜 그 점수를 받았는지 상세히 분석합니다.")

    col_deep_1, col_deep_2 = st.columns([2, 1])
    with col_deep_1:
        deep_query = st.text_input("분석할 검색어", value="냉각수펌프 전기실", key="deep_query_app")
    with col_deep_2:
        # Reuse target_filename as default if possible, otherwise generic default
        default_deep_file = ""
        deep_file = st.text_input("대상 파일 (필터 - 선택사항)", value=default_deep_file, key="deep_file_input_app")

    if st.button("🔬 랭킹 분석 실행", type="primary", use_container_width=True):
        st.markdown("### 1. 쿼리 분석 (Query Analysis)")
        
        # 1. Sanitization Logic
        import re
        sanitized_query = re.sub(r'\bAND\b', ' ', deep_query, flags=re.IGNORECASE)
        sanitized_query = re.sub(r'[&+\-|!(){}\[\]^"~*?:\\]', ' ', sanitized_query)
        sanitized_query = " ".join(sanitized_query.split())
        
        st.code(f"Original: '{deep_query}'\nSanitized: '{sanitized_query}'", language="text")
        
        keywords = sanitized_query.split()
        st.write(f"**Keywords extracted:** {keywords}")
        
        # 2. Execute Search
        st.markdown("### 2. 검색 결과 랭킹 (Top 20)")
        
        filter_expr = None
        if deep_file and deep_file.strip():
            filter_expr = f"search.ismatch('{deep_file}', 'metadata_storage_name')"
            
        with st.spinner("랭킹 분석 중..."):
            # Ensure search_manager is available (it's initialized at top level)
            manager = get_search_manager()
            results = manager.search(
                sanitized_query,
                filter_expr=filter_expr,
                search_mode="all", # Strict mode
                top=20,
                use_semantic_ranker=False # Raw score analysis
            )
            
        if not results:
            st.warning("검색 결과가 없습니다.")
        else:
            rank_data = []
            for i, doc in enumerate(results, 1):
                name = doc.get('metadata_storage_name', 'Unknown')
                content = doc.get('content', '')
                score = doc.get('@search.score', 0)
                
                # Keyword Matching Analysis
                content_upper = content.upper()
                matched_kws = []
                for kw in keywords:
                    if kw.upper() in content_upper:
                        matched_kws.append(kw)
                
                match_status = "✅ All" if len(matched_kws) == len(keywords) else f"⚠️ {len(matched_kws)}/{len(keywords)}"
                
                # Highlight specific pages
                highlight = ""
                if "(p.17)" in name: highlight = "🔴 (Issue)"
                if "(p.82)" in name: highlight = "🟢 (Target)"
                
                rank_data.append({
                    "Rank": i,
                    "Score": f"{score:.4f}",
                    "Page": f"{name} {highlight}",
                    "Match": match_status,
                    "Matched Keywords": ", ".join(matched_kws),
                    "Snippet": content[:100].replace("\n", " ") + "..."
                })
                
            st.dataframe(pd.DataFrame(rank_data), use_container_width=True)
            
            # Detailed Comparison
            st.markdown("### 3. 주요 페이지 상세 비교")
            target_pages = [d for d in results if "(p.82)" in d.get('metadata_storage_name', '') or "(p.17)" in d.get('metadata_storage_name', '')]
            
            if target_pages:
                for doc in target_pages:
                    name = doc.get('metadata_storage_name')
                    score = doc.get('@search.score')
                    st.markdown(f"#### 📄 {name} (Score: {score:.4f})")
                    st.text_area(f"Content of {name}", doc.get('content', ''), height=200)
            else:
                st.info("비교할 주요 페이지(p.17, p.82)가 Top 20 내에 없습니다.")
